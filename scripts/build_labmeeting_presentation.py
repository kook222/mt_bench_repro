#!/usr/bin/env python3
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
BASE_FIG = FIG / "base_paper"
OUT_DIR = ROOT / "presentation"
OUT_DIR.mkdir(exist_ok=True)
PPTX_PATH = OUT_DIR / "mt_bench_labmeeting_study_seminar_v2.pptx"
NOTES_MD_PATH = OUT_DIR / "mt_bench_labmeeting_study_seminar_v2_notes.md"


BG = RGBColor(250, 248, 244)
NAVY = RGBColor(19, 40, 58)
SLATE = RGBColor(54, 69, 79)
RUST = RGBColor(180, 92, 69)
SAGE = RGBColor(96, 122, 96)
GOLD = RGBColor(192, 148, 74)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(109, 119, 128)
LIGHT = RGBColor(233, 229, 219)
MIST = RGBColor(243, 239, 232)

FONT = "Apple SD Gothic Neo"
TITLE_SIZE = 30
BODY_SIZE = 18
SMALL_SIZE = 10
HEADER_TITLE_SIZE = 24


@dataclass
class SlideSpec:
    title: str
    section: str
    purpose: str
    bullets: list[str]
    notes: str
    rq_tag: str | None = None
    layout: str = "bullets"
    images: list[Path] = field(default_factory=list)
    image_captions: list[str] = field(default_factory=list)
    side_quote: str | None = None
    stat_boxes: list[tuple[str, str]] = field(default_factory=list)


def add_textbox(slide, left, top, width, height, text="", *, font_size=BODY_SIZE,
                bold=False, color=SLATE, align=PP_ALIGN.LEFT, fill=None,
                line=None, margin=0.08, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line
    else:
        box.line.fill.background()
    return box


def add_bullets(box, bullets: Iterable[str], *, font_size=BODY_SIZE, color=SLATE):
    tf = box.text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {bullet}"
        p.level = 0
        p.space_after = Pt(9)
        p.space_before = Pt(0)
        p.line_spacing = 1.15
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box


def add_header(slide, title: str, section: str, slide_no: int, total: int, rq_tag: str | None = None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0, 0, Inches(13.333), Inches(0.36)
    ).fill.solid()
    bar = slide.shapes[-1]
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_textbox(slide, Inches(0.62), Inches(0.5), Inches(10.9), Inches(0.62),
                title, font_size=HEADER_TITLE_SIZE, bold=True, color=NAVY)
    add_textbox(slide, Inches(10.55), Inches(0.05), Inches(1.7), Inches(0.22),
                section.upper(), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    if rq_tag:
        add_textbox(slide, Inches(9.2), Inches(0.05), Inches(1.1), Inches(0.22),
                    rq_tag, font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                    fill=LIGHT, line=LIGHT, margin=0.02, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(12.2), Inches(0.05), Inches(0.6), Inches(0.22),
                f"{slide_no}/{total}", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def add_footer(slide):
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.42), Inches(7.12), Inches(12.45), Inches(0.04)
    ).fill.solid()
    line = slide.shapes[-1]
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()
    add_textbox(
        slide, Inches(0.45), Inches(7.12), Inches(9.5), Inches(0.18),
        "MT-Bench base paper study + reproduction lab meeting", font_size=9, color=GRAY
    )


def add_purpose_box(slide, purpose: str):
    add_textbox(
        slide,
        Inches(0.55), Inches(0.78), Inches(12.2), Inches(0.52),
        f"이 슬라이드의 목적: {purpose}",
        font_size=13, bold=True, color=NAVY, fill=LIGHT, line=LIGHT, valign=MSO_ANCHOR.MIDDLE
    )


def add_image(slide, path: Path, left, top, width=None, height=None):
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_stat_box(slide, left, top, label, value, accent):
    add_textbox(slide, left, top, Inches(1.8), Inches(0.95), "",
                fill=WHITE, line=LIGHT, margin=0.12)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.12), top + Inches(0.12), Inches(0.48), Inches(0.22)
    )
    cap = slide.shapes[-1]
    cap.fill.solid()
    cap.fill.fore_color.rgb = accent
    cap.line.fill.background()
    add_textbox(slide, left + Inches(0.12), top + Inches(0.18), Inches(1.55), Inches(0.18),
                label, font_size=10, bold=True, color=WHITE, margin=0.02, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.42), Inches(1.56), Inches(0.35),
                value, font_size=21, bold=True, color=NAVY, margin=0.02)


def build_title_slide(prs: Presentation, spec: SlideSpec, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # background blocks
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    top = slide.shapes[-1]
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.9), Inches(0.9), Inches(4.433), Inches(6.6))
    side = slide.shapes[-1]
    side.fill.solid()
    side.fill.fore_color.rgb = MIST
    side.line.fill.background()

    add_textbox(slide, Inches(0.65), Inches(1.25), Inches(7.7), Inches(1.5),
                "오픈소스 LLM-as-a-Judge의\n신뢰도는 어디까지 믿을 수 있는가", font_size=31, bold=True, color=NAVY)
    add_textbox(slide, Inches(0.68), Inches(2.9), Inches(7.2), Inches(0.42),
                "MT-Bench 논문 스터디 + 재현 실험 + Git walkthrough", font_size=17, bold=True, color=RUST)
    add_textbox(slide, Inches(0.68), Inches(3.45), Inches(7.55), Inches(1.0),
                "원 논문을 먼저 해설한 뒤, 제가 같은 프로토콜을 어떻게 오픈소스 judge 실험으로 확장했는지\nPhase 1–6 결과와 저장소 구조까지 같이 설명하는 랩미팅 발표입니다.",
                font_size=18, color=SLATE)

    card_specs = [
        ("메인 judge", "Qwen 7B · 14B · 32B", RUST),
        ("보조 judge", "InternLM20B · GPT-4o-mini", SAGE),
        ("hold-out", "330 splits · 66K random draws", GOLD),
    ]
    for i, (label, value, accent) in enumerate(card_specs):
        box = add_textbox(slide, Inches(0.72 + i * 2.45), Inches(5.25), Inches(2.18), Inches(1.05), "",
                          fill=WHITE, line=LIGHT, margin=0.11)
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                               Inches(0.84 + i * 2.45), Inches(5.36), Inches(0.72), Inches(0.22))
        chip = slide.shapes[-1]
        chip.fill.solid()
        chip.fill.fore_color.rgb = accent
        chip.line.fill.background()
        add_textbox(slide, Inches(0.88 + i * 2.45), Inches(5.39), Inches(1.05), Inches(0.18),
                    label, font_size=9, bold=True, color=WHITE, margin=0.01, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(0.84 + i * 2.45), Inches(5.68), Inches(1.88), Inches(0.45),
                    value, font_size=15, bold=True, color=NAVY, margin=0.02)

    add_textbox(slide, Inches(9.1), Inches(1.3), Inches(3.7), Inches(0.35),
                "오늘 발표에서 답할 질문", font_size=17, bold=True, color=NAVY)
    q_box = add_textbox(slide, Inches(9.1), Inches(1.78), Inches(3.6), Inches(4.0), "",
                        fill=WHITE, line=LIGHT)
    add_bullets(q_box, [
        "베이스 논문은 무엇을 제안했고 왜 중요했는가?",
        "내 재현 코드는 어디까지 원문을 따라갔고 어디서 확장했는가?",
        "오픈소스 judge를 실제로 믿어도 되는가?",
        "비용을 줄이면서도 순위를 얼마나 보존할 수 있는가?",
    ], font_size=16)
    add_textbox(slide, Inches(9.1), Inches(6.12), Inches(3.7), Inches(0.55),
                "박승현\nCLINK Lab, Pusan National University", font_size=16, color=SLATE)
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = spec.notes
    return slide


def build_bullets_slide(prs: Presentation, spec: SlideSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title, spec.section, slide_no, total, spec.rq_tag)
    if spec.side_quote:
        add_textbox(slide, Inches(0.7), Inches(1.05), Inches(8.8), Inches(0.26),
                    spec.side_quote, font_size=13, bold=True, color=RUST)

    if spec.images:
        fig_top = Inches(1.42)
        if len(spec.images) == 1:
            add_image(slide, spec.images[0], Inches(0.62), fig_top, width=Inches(8.18))
            cap = spec.image_captions[0] if spec.image_captions else spec.images[0].name
            add_textbox(slide, Inches(0.75), Inches(6.0), Inches(7.9), Inches(0.26),
                        cap, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
        else:
            add_image(slide, spec.images[0], Inches(0.62), fig_top, width=Inches(3.95))
            add_image(slide, spec.images[1], Inches(4.82), fig_top, width=Inches(3.95))
            if spec.image_captions:
                add_textbox(slide, Inches(0.7), Inches(5.06), Inches(3.8), Inches(0.24),
                            spec.image_captions[0], font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
                add_textbox(slide, Inches(4.9), Inches(5.06), Inches(3.8), Inches(0.24),
                            spec.image_captions[1], font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

        body_box = add_textbox(slide, Inches(8.95), Inches(1.42), Inches(3.55), Inches(5.15), "",
                               fill=WHITE, line=LIGHT, margin=0.14)
        add_bullets(body_box, spec.bullets, font_size=17)
        if spec.stat_boxes:
            for i, (label, value) in enumerate(spec.stat_boxes):
                accent = [RUST, SAGE, GOLD, NAVY][i % 4]
                x = 8.9 + (i % 2) * 1.82
                y = 5.92 + (i // 2) * 0.9
                add_stat_box(slide, Inches(x), Inches(y), label, value, accent)
    else:
        body_box = add_textbox(slide, Inches(0.68), Inches(1.42), Inches(12.0), Inches(5.45), "",
                               fill=WHITE, line=LIGHT, margin=0.16)
        add_bullets(body_box, spec.bullets, font_size=20)
        if spec.stat_boxes:
            for i, (label, value) in enumerate(spec.stat_boxes):
                accent = [RUST, SAGE, GOLD, NAVY][i % 4]
                x = 0.9 + i * 2.95
                add_stat_box(slide, Inches(x), Inches(5.92), label, value, accent)

    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = spec.notes
    return slide


def build_repo_slide(prs: Presentation, spec: SlideSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title, spec.section, slide_no, total, spec.rq_tag)

    add_textbox(slide, Inches(0.65), Inches(1.05), Inches(7.2), Inches(0.26),
                "발표를 보고 저장소를 열었을 때 바로 검증 가능한 구조가 되도록 정리했습니다.", font_size=13, bold=True, color=RUST)

    add_textbox(slide, Inches(0.65), Inches(1.35), Inches(4.1), Inches(0.3),
                "핵심 설계 결정", font_size=18, bold=True, color=NAVY)
    tree = add_textbox(slide, Inches(0.65), Inches(1.72), Inches(4.0), Inches(4.95), "",
                       fill=WHITE, line=LIGHT, margin=0.12)
    add_bullets(tree, [
        "JSONL append + resume로 중간 실패가 나도 완료분을 보존",
        "pairwise는 AB/BA 두 방향을 모두 실행하고 불일치는 inconsistent로 처리",
        "generation temperature 0.7, judge temperature 0.0으로 역할을 분리",
        "raw judgment → aggregate CSV → figure → paper까지 한 체인으로 연결",
    ], font_size=16)

    add_textbox(slide, Inches(4.95), Inches(1.35), Inches(3.0), Inches(0.3),
                "Repo 핵심 구조", font_size=18, bold=True, color=NAVY)
    flow = add_textbox(slide, Inches(4.95), Inches(1.72), Inches(3.2), Inches(4.95), "",
                       fill=WHITE, line=LIGHT)
    add_bullets(flow, [
        "`src/mtbench_repro/` — client, prompts, aggregate, CLI",
        "`scripts/` — generation, judge, phase 분석, 발표용 재집계",
        "`data/` — answer / judgment / summary csv",
        "`figures/` — 논문 및 발표 그림",
        "`paper/`, `presentation/` — 원고와 발표 산출물",
    ], font_size=17)

    add_textbox(slide, Inches(8.45), Inches(1.35), Inches(4.1), Inches(0.3),
                "검증 루프", font_size=18, bold=True, color=NAVY)
    repro = add_textbox(slide, Inches(8.45), Inches(1.72), Inches(4.2), Inches(4.95), "",
                        fill=WHITE, line=LIGHT)
    add_bullets(repro, [
        "결론 문장 → figure → csv → raw jsonl로 역추적 가능",
        "phase별 raw judgment와 summary csv를 분리",
        "mock / real 결과 분리 및 reference 별도 집계",
        "발표 자료와 대본까지 script 기반으로 재생성 가능",
    ], font_size=16)
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = spec.notes
    return slide


def build_timeline_slide(prs: Presentation, spec: SlideSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title, spec.section, slide_no, total, spec.rq_tag)
    add_textbox(slide, Inches(0.65), Inches(1.05), Inches(7.8), Inches(0.26),
                "Phase는 시간순 나열이 아니라, 신뢰도를 단계적으로 끌어올리는 구조로 설계했습니다.", font_size=13, bold=True, color=RUST)

    phases = [
        ("P1", "Self-judge bias", RUST),
        ("P2", "예비 비교", GOLD),
        ("P3", "메인 scaling", NAVY),
        ("P4", "InternLM", SAGE),
        ("P5", "GPT-4o-mini", RUST),
        ("P6", "330 split hold-out", NAVY),
    ]
    y = Inches(1.9)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), y+Inches(0.42), Inches(11.5), Inches(0.06))
    line = slide.shapes[-1]
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()

    x_positions = [1.0, 2.9, 4.8, 6.7, 8.6, 10.5]
    for x, (tag, label, accent) in zip(x_positions, phases):
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), y, Inches(0.72), Inches(0.72))
        circ = slide.shapes[-1]
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent
        circ.line.fill.background()
        add_textbox(slide, Inches(x), y+Inches(0.13), Inches(0.72), Inches(0.18),
                    tag, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0.01)
        add_textbox(slide, Inches(x-0.22), y+Inches(0.88), Inches(1.2), Inches(0.5),
                    label, font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    panel = add_textbox(slide, Inches(0.85), Inches(3.25), Inches(11.65), Inches(2.8), "",
                        fill=WHITE, line=LIGHT, margin=0.14)
    add_bullets(panel, spec.bullets, font_size=17)
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = spec.notes
    return slide


def build_summary_slide(prs: Presentation, spec: SlideSpec, slide_no: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title, spec.section, slide_no, total, spec.rq_tag)
    left = add_textbox(slide, Inches(0.65), Inches(1.42), Inches(6.15), Inches(4.95), "",
                       fill=WHITE, line=LIGHT)
    add_bullets(left, spec.bullets, font_size=18)
    quote = spec.side_quote or "핵심 결론"
    add_textbox(slide, Inches(7.05), Inches(1.42), Inches(5.15), Inches(0.72),
                quote, font_size=19, bold=True, color=WHITE, fill=NAVY, line=NAVY, valign=MSO_ANCHOR.MIDDLE)
    if spec.stat_boxes:
        for i, (label, value) in enumerate(spec.stat_boxes):
            accent = [RUST, SAGE, GOLD, NAVY][i % 4]
            x = 7.12 + (i % 2) * 2.7
            y = 2.6 + (i // 2) * 1.28
            add_stat_box(slide, Inches(x), Inches(y), label, value, accent)
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = spec.notes
    return slide


def build_notes_md(slides: list[SlideSpec]):
    lines = ["# MT-Bench Labmeeting / Paper Study Deck Notes", ""]
    for idx, slide in enumerate(slides, start=1):
        lines.append(f"## Slide {idx}. {slide.title}")
        lines.append(f"- Section: {slide.section}")
        lines.append(f"- Purpose: {slide.purpose}")
        lines.append(f"- Layout: {slide.layout}")
        lines.append("- Bullets:")
        for bullet in slide.bullets:
            lines.append(f"  - {bullet}")
        if slide.images:
            lines.append("- Visuals:")
            for img, cap in zip(slide.images, slide.image_captions or [p.name for p in slide.images]):
                lines.append(f"  - {img.name}: {cap}")
        if slide.stat_boxes:
            lines.append("- Stat boxes:")
            for label, value in slide.stat_boxes:
                lines.append(f"  - {label}: {value}")
        lines.append("")
        lines.append("### Speaker Notes")
        lines.append(slide.notes.strip())
        lines.append("")
    NOTES_MD_PATH.write_text("\n".join(lines), encoding="utf-8")


def get_slide_specs() -> list[SlideSpec]:
    return [
        SlideSpec(
            title="오픈소스 LLM-as-a-Judge의 신뢰도 분석",
            section="Title",
            purpose="청중에게 오늘 발표의 범위와 정체성을 한 번에 전달한다.",
            bullets=[],
            notes="""안녕하세요. 오늘 발표는 두 가지를 동시에 하려는 자리입니다.
첫째는 베이스 논문인 MT-Bench/Chatbot Arena 논문이 정확히 무엇을 제안했는지 다시 이해하는 논문 스터디이고,
둘째는 제가 그 위에서 어떤 재현과 확장을 했는지를 보여드리는 랩미팅입니다.

그래서 발표 구조도 의도적으로 두 층으로 나눴습니다. 먼저 원 논문이 왜 등장했는지,
MT-Bench와 LLM-as-a-Judge가 어떤 문제를 풀려고 했는지 설명한 다음,
그다음부터는 제 Git 저장소 기준으로 phase별 실험을 따라가겠습니다.

오늘의 핵심 메시지는 세 가지입니다.
첫째, 오픈소스 judge도 크기를 키우면 꽤 일관된 평가자가 될 수 있습니다.
둘째, 다만 inconsistency가 줄어드는 것과 공정성이 높아지는 것은 같은 말이 아닙니다.
셋째, 문항 수를 줄이는 문제는 same-set에서는 40문항이 강하지만,
반복 hold-out으로 보면 60문항이 더 안전한 운영점이라는 것입니다.""",
        ),
        SlideSpec(
            title="오늘 발표의 큰 흐름",
            section="Roadmap",
            purpose="20분 동안 무엇을 어떤 순서로 듣게 되는지 먼저 안정적으로 안내한다.",
            bullets=[
                "1부: 베이스 논문이 해결하려는 문제와 핵심 기여",
                "2부: 내 재현 파이프라인과 Phase 1–6 실험 구조",
                "3부: 메인 결과, Git 정리 방식, 그리고 최종 해석",
            ],
            notes="""이 슬라이드는 사실상 청중의 집중을 정렬하는 슬라이드입니다.
교수님이나 심사위원이 가장 불편해하는 건 발표자가 논문 설명과 자기 결과를 뒤섞어버리는 경우입니다.
그래서 저는 먼저 베이스 논문을 독립적으로 설명하고, 그다음에 제 실험을 얹는 구조로 가겠습니다.

1부에서는 원 논문이 왜 기존 벤치마크를 문제 삼았는지, MT-Bench와 Chatbot Arena를 왜 만들었는지,
그리고 GPT-4 judge를 어떻게 검증했는지를 짚겠습니다.
2부에서는 제 Git 저장소 관점에서 파이프라인을 설명합니다. 어떤 코드가 있고, 어떤 phase가 있고,
어디까지가 재현이고 어디부터가 확장인지 분리해서 말씀드리겠습니다.
3부에서는 결과 해석과 결론입니다. 특히 same-set 결과와 repeated hold-out 결과를 분리해서 읽는 법,
그리고 제가 이 연구를 어디까지 주장하고 어디서부터는 주장하지 않는지까지 솔직하게 정리하겠습니다.""",
        ),
        SlideSpec(
            title="베이스 논문의 문제의식과 핵심 자산",
            section="Base Paper",
            purpose="원 논문이 왜 등장했고, 어떤 benchmark와 human evidence를 만들었는지 한 장에서 잡는다.",
            bullets=[
                "기존 객관식 벤치마크는 open-ended, multi-turn 선호를 충분히 반영하지 못함",
                "사람 선호를 직접 모으는 평가는 느리고 비싸서 반복 실험이 어려움",
                "MT-Bench 80문항과 Chatbot Arena를 함께 만들어 controlled benchmark와 in-the-wild preference를 동시에 확보",
            ],
            side_quote="Base paper question: “Can a strong LLM approximate human preference at scale?”",
            images=[BASE_FIG / "paper_chatbot_arena_ui.png", BASE_FIG / "paper_mtbench_winrate_fig3.png"],
            image_captions=["Original paper: Chatbot Arena UI screenshot (Figure 19)", "Original paper: MT-Bench average win-rate curves (Figure 3)"],
            notes="""원 논문은 아주 명확한 문제의식에서 출발합니다.
기존 LLM 벤치마크, 예를 들어 MMLU나 HELM류는 모델의 핵심 능력을 보긴 하지만,
사용자가 실제로 좋아하는 대화형 응답 품질과는 어긋나는 경우가 많다는 겁니다.

예를 들어 instruction following, multi-turn consistency, 답변의 유용성 같은 건
정답이 하나로 떨어지는 객관식 점수만으로는 잘 측정되지 않습니다.
그렇다고 사람 평가만 계속 하자니 비용과 시간이 너무 큽니다.

그래서 이 논문은 강한 LLM, 당시에는 GPT-4를 judge로 써서
인간 선호를 근사할 수 있는지 체계적으로 보자고 합니다.
즉 질문 자체가 “LLM이 다른 LLM을 채점해도 되나?”가 아니라,
“그 판단이 인간 선호와 어느 정도까지 맞느냐”에 가깝습니다.

이때 만든 핵심 자산이 MT-Bench와 Chatbot Arena입니다.
MT-Bench는 8개 카테고리의 80문항 multi-turn benchmark이고,
Arena는 실제 사용자가 익명으로 두 모델 중 더 나은 답변에 투표하는 플랫폼입니다.
즉 원 논문은 benchmark 하나를 제안한 게 아니라,
controlled benchmark와 wild preference data를 동시에 만든 논문입니다.""",
        ),
        SlideSpec(
            title="베이스 논문의 judge 프로토콜과 핵심 수치",
            section="Base Paper",
            purpose="우리가 따라갈 프로토콜과 원 논문의 기준 수치를 한 장에 묶는다.",
            bullets=[
                "GPT-4 judge는 position, verbosity, self-enhancement bias 가능성을 함께 분석",
                "MT-Bench에서 GPT-4 judge와 인간 expert의 non-tie agreement는 85%, 인간-인간은 81%",
                "Chatbot Arena에서도 GPT-4와 crowd human의 non-tie agreement는 87%, GPT-4 single은 85%",
                "MT-Bench score는 LLaMA-13B 2.61, Vicuna-13B 6.39, GPT-3.5 7.94, GPT-4 8.99",
            ],
            side_quote="원 논문의 결론은 ‘GPT-4 judge가 완벽하다’가 아니라, ‘인간 선호를 실용적으로 근사할 수 있다’였습니다.",
            images=[BASE_FIG / "paper_mtbench_agreement_table5.png", BASE_FIG / "paper_table8_scores.png"],
            image_captions=["Original paper: MT-Bench judge-human agreement table (Table 5)", "Original paper: MT-Bench score table for model variants (Table 8)"],
            notes="""이 슬라이드는 원 논문에서 우리가 기억해야 할 프로토콜과 핵심 수치를 한 번에 정리하는 슬라이드입니다.
발표에서 여기까지가 사실상 논문 스터디 파트의 핵심입니다.

첫째, MT-Bench 점수와 Arena human preference가 높은 상관을 보였다는 점.
이게 원 논문이 영향력을 얻은 가장 큰 이유입니다.
둘째, single-answer grading과 pairwise ranking이 대체로 수렴한다는 점.
셋째, hard category에서 모델 간 차이가 더 크게 드러난다는 점입니다.

그런데 이 세 결론을 곧바로 ‘judge는 믿어도 된다’로 읽으면 안 됩니다.
논문도 bias를 아주 분명히 분석했고, 특히 position bias와 verbosity bias를 숨기지 않았습니다.
즉 원 논문은 perfect judge를 선언한 게 아니라, usable surrogate를 제안한 겁니다.

이 발표에서 제 실험은 바로 이 usable surrogate라는 표현을 오픈소스 judge 맥락에서 다시 묻는 작업입니다.
즉, GPT-4에서는 usable했는데 Qwen/InternLM/GPT-4o-mini의 조합에서는 어디까지 usable한가를 보는 겁니다.""",
        ),
        SlideSpec(
            title="무엇을 그대로 가져오고, 무엇을 바꾸며, 어떤 RQ를 세웠는가",
            section="Base Paper",
            purpose="재현과 확장의 경계를 분명히 하고 이후 Phase를 RQ와 연결한다.",
            rq_tag="RQ1-4",
            bullets=[
                "그대로 가져온 것: 80문항 MT-Bench, single/pairwise/reference 3종 채점 구조",
                "바꾼 것: GPT-4 judge를 Qwen·InternLM·GPT-4o-mini로 대체하고 raw artifact를 전부 저장",
                "RQ1: judge 크기와 reliability는 어떻게 바뀌는가?",
                "RQ2: 더 큰 judge에서 남는 residual bias는 무엇인가?",
                "RQ3: cross-family / external judge에서도 같은 broad rank가 유지되는가?",
                "RQ4: 문항 수를 줄여도 hold-out 모델에서 서열을 보존할 수 있는가?",
            ],
            side_quote="이 발표는 ‘원 논문을 흉내 낸 실험’이 아니라, ‘원 논문의 프로토콜을 재현 가능한 오픈소스 파이프라인으로 옮긴 작업’입니다.",
            notes="""이 슬라이드는 베이스 논문 설명과 제 실험 설명 사이의 브리지입니다.
청중 입장에서는 여기서 ‘그래서 도대체 어디까지가 재현이고, 어디부터가 새로운 분석이냐’를 알고 싶어합니다.

그대로 가져온 요소는 명확합니다.
MT-Bench 80문항, single-grade, pairwise, reference-guided grading이라는 세 가지 judge 프로토콜,
그리고 conservative한 pairwise 해석 철학은 전부 원 논문에서 그대로 이어받았습니다.

바꾼 부분도 분명합니다.
judge를 GPT-4 하나로 고정하지 않고, 오픈소스 Qwen/InternLM과 외부 API judge인 GPT-4o-mini까지 붙였습니다.
그리고 결과를 논문 표 수준이 아니라 raw JSONL, aggregated CSV, figure까지 모두 저장해서
다시 계산 가능한 형태로 만들었습니다.

마지막으로 확장한 질문이 있습니다.
judge scaling, residual position bias, abstain ensemble, tinyMT-Bench, exhaustive hold-out은
원 논문이 직접 답하지 않았던 질문들입니다.
즉 제 작업은 재현으로 시작하지만, 발표의 후반부는 재현 결과 위에 얹은 실증 확장이라고 보시면 됩니다.""",
            layout="summary",
            stat_boxes=[("재현", "프로토콜"), ("대체", "오픈 judge"), ("확장", "RQ1–RQ4"), ("산출물", "raw→paper")],
        ),
        SlideSpec(
            title="내 Git 저장소는 무엇을 재현 가능하게 만들었나",
            section="Repo",
            purpose="코드베이스가 단순 결과 저장소가 아니라 재현 가능한 연구 파이프라인이라는 점을 보여준다.",
            bullets=[
                "append+resume, conservative verdict, temperature 분리 같은 판단을 코드 레벨에 고정",
                "raw judgment와 aggregate CSV, figure, paper를 같은 체인으로 다시 생성 가능",
                "즉 결과를 보여주는 저장소가 아니라 결과를 다시 만드는 저장소로 설계",
            ],
            notes="""이 슬라이드는 결과보다도 재현성에 관한 슬라이드입니다.
심사나 랩미팅에서 자주 받는 질문이 “그래서 이건 한 번 돌린 실험인가, 아니면 다시 재현 가능한가?”입니다.

제 저장소는 그 점을 의식해서 구성했습니다.
src 안에는 judge client, prompt, aggregation, CLI가 있고,
scripts 안에는 generation, judge 실행, phase별 분석, 발표용 재집계가 있습니다.
data는 raw answer와 judgment, summary csv를 분리해 놓았고,
figures는 논문과 발표에 직접 들어가는 그림을 다시 생성할 수 있게 관리했습니다.

즉 이 repo는 결과를 예쁘게 모아둔 폴더가 아니라,
질문 생성부터 평가, 집계, figure, 원고까지 연결되는 end-to-end 실험 환경입니다.
나중에 마지막 슬라이드에서 다시 보겠지만,
이 구조 덕분에 Phase 6 hold-out도 추가 GPU 없이 재분석만으로 강화할 수 있었습니다.""",
            layout="repo",
        ),
        SlideSpec(
            title="Git에서 실제로 어디를 보면 무엇이 보이는가",
            section="Repo",
            purpose="랩미팅 이후에도 저장소를 읽을 수 있도록 산출물 탐색 경로를 구체적으로 알려준다.",
            bullets=[
                "`README.md`는 전체 서사, `data/`는 raw judgment와 요약 CSV, `figures/`는 발표/논문 그림",
                "`paper/`에는 blind/non-blind 원고와 PDF, `presentation/`에는 오늘 발표자료와 대본",
                "질문이 생기면 ‘결론 → figure → csv → raw jsonl’ 순서로 역추적하면 됨",
            ],
            side_quote="교수님이 저장소를 훑으실 때는 README만 보지 말고, 결과 figure 옆의 CSV와 raw judgment까지 바로 연결되도록 설계했습니다.",
            notes="""이 슬라이드는 발표 이후를 위한 안내서입니다.
랩미팅이 끝나면 보통 교수님이나 동료가 저장소를 직접 열어보게 되는데,
그때 어디부터 봐야 하는지 모르면 좋은 repo도 금방 피곤해집니다.

그래서 저는 산출물을 읽는 경로를 일부러 단순하게 만들었습니다.
README는 서사와 핵심 메시지, data는 raw answer와 judgment, summary CSV,
figures는 발표와 논문에 들어가는 시각화, paper는 원고와 PDF,
presentation은 오늘 발표 파일과 노트를 담습니다.

실제로 검증이 필요할 때는 결론 문장 하나에서 출발해서
해당 figure, 그 figure를 만든 CSV, 마지막으로 raw judgment JSONL까지 역으로 따라가면 됩니다.
이건 단순 정리 습관이 아니라 연구 신뢰도를 높이는 설계라고 생각합니다.

발표에서 이 슬라이드를 짧게라도 넣는 이유는,
제 작업이 단순히 ‘결과가 좋다’가 아니라 ‘다시 확인 가능하다’는 점을 보여주기 위해서입니다.""",
            layout="summary",
            stat_boxes=[("README", "서사"), ("data", "raw + csv"), ("figures", "발표/논문"), ("presentation", "ppt + notes")],
        ),
        SlideSpec(
            title="Phase 1–6 전체 실험 흐름",
            section="Repo",
            purpose="발표 후반 결과를 이해할 수 있도록 phase별 역할을 미리 정렬한다.",
            bullets=[
                "Phase 1–2: self-judge 편향과 예비 외부 judge sanity check",
                "Phase 3: 메인 same-family scaling 실험(Qwen 7B/14B/32B)",
                "Phase 4–5: InternLM20B, GPT-4o-mini로 cross-family / external check",
                "Phase 6: 11개 모델 풀에서 exhaustive leave-4-out hold-out",
            ],
            notes="""이 슬라이드는 발표 전체의 지도 역할을 합니다.
Phase를 숫자로만 나열하면 지루해지기 쉬운데, 각 phase의 존재 이유를 기억해 두시면 뒤 슬라이드가 훨씬 잘 들어옵니다.

Phase 1은 self-judge bias를 직접 확인하는 단계입니다.
Phase 2는 외부 14B judge로 예비 sanity check를 하는 단계고,
여기까지는 말 그대로 파이프라인이 제대로 도는지, 큰 그림이 맞는지 보는 준비운동입니다.

진짜 메인은 Phase 3입니다. Qwen 7B, 14B, 32B를 같은 family 안에서 비교하면서
judge scaling을 관찰합니다. Phase 4와 5는 여기서 나온 결과가 Qwen 고유 현상이 아닌지 보는 보조 검증입니다.
마지막 Phase 6은 문항 축소 결과를 hold-out 모델에서도 다시 보는 단계입니다.

이 발표에서는 Phase 3를 중심축으로 두고,
Phase 4~6이 그 주장을 어디까지 보강하고 어디서부터는 한계를 드러내는지를 보여드리겠습니다.""",
            layout="timeline",
        ),
        SlideSpec(
            title="Phase 1–2: self-judge 편향과 예비 sanity check",
            section="Results",
            purpose="외부 judge가 왜 필요한지, 그리고 메인 실험 전 무엇을 확인했는지 보여준다.",
            rq_tag="Prelude",
            bullets=[
                "Qwen2.5-7B self-judge는 overall 8.12인데 Math·Coding은 각각 8.80으로 튀어 오름",
                "외부 14B judge 도입 후 서열은 안정화되지만 pairwise inconsistency는 여전히 높음",
                "즉 ‘오픈소스 judge를 쓰자’보다 ‘어떤 judge를 어떻게 쓰나’가 더 중요함",
            ],
            images=[FIG / "fig0_phase1_scores.png", FIG / "fig2_overall_rankings.png"],
            image_captions=["bar chart: Phase 1 self-judge category scores", "ranking plot: 초기 overall ordering"],
            notes="""Phase 1과 2는 메인 결과로 쓰기보다는,
왜 이후의 실험이 필요한지를 보여주는 출발점입니다.

먼저 self-judge를 해보면 Qwen2.5-7B가 자기 자신이나 유사 강점이 있는 카테고리를 높게 주는 경향이 드러납니다.
특히 Math, Coding 쪽이 전형적입니다.
이건 LLM-as-a-Judge를 쓸 때 가장 먼저 피해야 할 함정 중 하나입니다.

그래서 외부 judge를 도입한 게 Phase 2입니다.
그런데 외부 judge를 쓰는 것만으로 문제가 끝나지는 않습니다.
서열은 어느 정도 정리되지만, pairwise inconsistency 자체는 여전히 큽니다.
즉 “self-judge만 아니면 된다”가 아니라,
judge의 크기, 아키텍처, prompt 안정성까지 같이 봐야 한다는 교훈을 줍니다.

이 슬라이드의 결론은 단순합니다.
Phase 1–2는 메인 claim이 아니라, 메인 실험을 설계하게 만든 문제 제기 단계입니다.""",
        ),
        SlideSpec(
            title="Phase 3: Qwen same-family scaling이 보여준 메인 결과",
            section="Results",
            purpose="Judge 크기 증가가 inconsistency를 어떻게 바꾸는지 메인 empirical trend를 각인시킨다.",
            rq_tag="RQ1",
            bullets=[
                "Pairwise inconsistency: 7B 78.75% → 14B 46.85% → 32B 32.86%",
                "점수 범위도 0.84pt → 1.13pt → 1.47pt로 확대되어 변별력 증가",
                "메인 claim은 `Qwen same-family empirical trend`로 제한해서 읽어야 함",
            ],
            images=[FIG / "fig4_judge_scaling.png"],
            image_captions=["bar chart: judge scaling과 카테고리별 inconsistency"],
            notes="""이 슬라이드가 발표의 중심입니다.
제가 가장 강하게 주장할 수 있는 건 여기서 나온 same-family scaling 결과입니다.

핵심 수치는 아주 단순합니다.
Qwen judge를 7B에서 14B, 다시 32B로 키우면 pairwise inconsistency가 78.75%에서 32.86%까지 단조 감소합니다.
이건 굉장히 큰 변화입니다.
동시에 single-grade 점수 범위도 더 넓어지기 때문에 모델 간 차이를 구분하는 능력도 좋아집니다.

다만 여기서 중요한 건 주장 수위를 지키는 겁니다.
이건 `모든 judge에서 스케일링 법칙이 성립한다`가 아닙니다.
정확히는 `Qwen2.5 동일 패밀리에서는 judge가 커질수록 inconsistency가 줄었다`는 empirical trend입니다.

저는 발표에서도 이 표현을 일부러 보수적으로 쓰고 있습니다.
왜냐하면 이 한 슬라이드가 아무리 강해도, 같은 family의 세 점만 갖고 universal law처럼 말하면 바로 공격 포인트가 되기 때문입니다.
그래서 뒤의 Phase 4와 5는 이 메인 슬라이드를 보조하는 역할로 보시면 됩니다.""",
        ),
        SlideSpec(
            title="Phase 3 추가 확인: single-grade 서열은 얼마나 안정적인가",
            section="Results",
            purpose="pairwise reliability와 별개로 single-grade 서열 자체가 얼마나 유지되는지 보여준다.",
            rq_tag="RQ1",
            bullets=[
                "32B single-grade 기준 상위권은 gemma·Phi, 하위권은 SOLAR·Zephyr로 정리됨",
                "Qwen judge 간 Spearman ρ는 0.750–0.821로 broad ranking이 유지되며, 7B↔14B bootstrap 95% CI는 [0.643, 0.964]",
                "즉 judge 크기 차이는 ‘서열 자체’보다 ‘pairwise reliability’에서 더 크게 드러남",
            ],
            images=[FIG / "fig5_phase3_scores.png"],
            image_captions=["Phase 3 single-grade 점수 분포"],
            stat_boxes=[("ρ 7↔14", "0.821"), ("ρ 7↔32", "0.786"), ("ρ 14↔32", "0.750")],
            notes="""이 슬라이드는 pairwise 결과와 single-grade 결과를 분리해서 읽게 해주는 슬라이드입니다.
앞 슬라이드에서는 judge가 커질수록 inconsistency가 줄어든다는 점을 봤고,
이 슬라이드에서는 그러면 모델 서열 자체는 얼마나 흔들리는지를 따로 봅니다.

핵심은 생각보다 single-grade rank는 꽤 안정적이라는 점입니다.
judge 7B, 14B, 32B를 바꿔도 상위권과 하위권의 큰 흐름은 유지됩니다.
즉 작은 judge가 완전히 쓸모없는 게 아니라, broad rank sanity check 정도는 할 수 있다는 뜻입니다.

문제는 거기서 멈추면 안 된다는 겁니다.
single-grade broad rank가 유지된다고 해서 pairwise verdict 품질까지 좋은 것은 아닙니다.
그래서 제가 이 발표 내내 single-grade rank consistency와 pairwise decision reliability를 분리해서 말하고 있습니다.

이 슬라이드는 그 구분을 청중에게 명확히 남기는 역할을 합니다.
그리고 뒤의 InternLM, GPT-4o-mini 슬라이드도 사실 이 프레임으로 읽는 게 맞습니다.""",
        ),
        SlideSpec(
            title="Phase 3 심화: residual bias와 문항 변별도",
            section="Results",
            purpose="단순히 inconsistency가 줄었다는 사실을 넘어, 남는 오류의 성격까지 설명한다.",
            rq_tag="RQ2",
            bullets=[
                "32B judge의 남은 불일치 중 94.9%가 first-position win으로 연결",
                "즉 judge가 커질수록 ‘남는 오류’는 위치 편향 중심으로 농축됨",
                "문항 변별도는 Math·Coding·Reasoning에서 특히 높아 tinyMT의 근거가 됨",
            ],
            images=[FIG / "fig11_position_bias.png", FIG / "fig8_discriminability.png"],
            image_captions=["bar chart: residual position bias", "bar chart: 문항 변별도 분포"],
            notes="""이 슬라이드는 개인적으로 가장 흥미로운 결과입니다.
많은 사람이 judge가 커지면 그냥 더 좋아진다고 생각하기 쉬운데,
제가 본 건 조금 다릅니다.

전체 inconsistency는 줄어듭니다. 그런데 남아 있는 불일치만 떼어 놓고 보면,
그 대부분이 first-position bias와 연결됩니다.
즉 작은 judge는 여기저기서 흔들리지만,
큰 judge는 흔들리는 경우 자체는 줄어들되 남는 흔들림이 순서 민감성으로 집중된다는 뜻입니다.

이건 해석상 중요한 포인트입니다.
‘더 큰 judge = 더 공정한 judge’가 아니라,
‘더 큰 judge = 더 일관된 judge, 하지만 남은 불일치의 구조는 더 편향적일 수 있음’입니다.

오른쪽의 문항 변별도 결과는 다음 슬라이드들의 bridge입니다.
Math, Coding, Reasoning에서 model score standard deviation이 크다는 건,
이 문항들이 모델 서열을 가르는 데 더 유용하다는 의미입니다.
즉 tinyMT-Bench는 갑자기 등장한 아이디어가 아니라,
이 슬라이드의 변별도 분석에서 자연스럽게 이어진 결과입니다.""",
        ),
        SlideSpec(
            title="앙상블은 정말 큰 judge를 이길 수 있는가",
            section="Results",
            purpose="실무적으로 가장 궁금한 질문 하나에 대해 명확한 yes/no를 준다.",
            rq_tag="RQ2",
            bullets=[
                "단순 다수결 앙상블은 오히려 단일 32B보다 나쁨: 58.63% > 32.86%",
                "inconsistent를 기권 처리한 abstain ensemble은 24.70%로 가장 낮음",
                "핵심 교훈: 저품질 judge를 그냥 합치면 평균이 아니라 오염이 됨",
            ],
            images=[FIG / "fig13_ensemble_v2.png"],
            image_captions=["bar chart: 다수결 vs 기권 앙상블 비교"],
            notes="""실무적으로는 이 슬라이드가 제일 즉답형입니다.
많은 사람이 작은 judge 여러 개를 합치면 큰 judge 하나보다 나아질 거라고 기대합니다.
그런데 제 결과는 그렇지 않다고 말합니다.

단순 다수결은 오히려 더 나빠집니다.
왜냐하면 7B의 noisy한 inconsistent 표가 앙상블에 그대로 유입되기 때문입니다.
이건 ensemble이 wisdom of crowds가 되려면 crowd 자체가 최소한의 품질을 가져야 한다는 걸 보여줍니다.

반대로 abstain 방식은 좋았습니다.
inconsistent를 억지로 표로 세지 않고 기권으로 처리하면,
결정적인 표가 있는 경우에만 winner를 선언하게 되고,
결국 inconsistency가 24.70%까지 내려갑니다.

이 결과는 학술적으로도 흥미롭지만,
실제로 judge system을 설계하는 입장에서는 더 유용합니다.
즉 ‘모델 여러 개를 붙일 거면 voting rule부터 다시 설계하라’는 교훈을 줍니다.
발표에서는 이 슬라이드를 practical takeaway로 강조하겠습니다.""",
        ),
        SlideSpec(
            title="Phase 4: InternLM2.5는 무엇을 보여주고 무엇은 못 보여주는가",
            section="Results",
            purpose="보조 cross-family judge의 가치와 실패 사례를 동시에 보여줘 신뢰도를 높인다.",
            rq_tag="RQ3",
            bullets=[
                "InternLM2.5-20B는 Qwen32와 Spearman ρ=0.893로 broad ranking을 유지",
                "반면 InternLM2.5-7B pairwise는 winner=error가 72.6%로 실전 judge로 부적합",
                "즉 ‘cross-family check는 됨’, 하지만 ‘아무 7B judge나 pairwise에 쓰면 안 됨’",
            ],
            images=[FIG / "fig17_phase4_internlm.png"],
            image_captions=["summary chart: Phase 4 InternLM judge 결과"],
            notes="""Phase 4는 일부러 좋은 결과와 나쁜 결과를 같이 보여주는 슬라이드입니다.
왜냐하면 이게 오히려 발표 전체의 신뢰도를 높이기 때문입니다.

InternLM2.5-20B는 생각보다 괜찮습니다.
Qwen32와 Spearman rho가 0.893이기 때문에,
적어도 seen-7 수준에서는 ‘모델 서열의 큰 흐름’이 유지된다고 말할 수 있습니다.

하지만 7B는 다릅니다.
InternLM2.5-7B는 pairwise verdict format을 안정적으로 따르지 못해서,
전체 record의 72.6%가 사실상 error로 남았습니다.
이건 단순 파서 버그가 아니라, 작은 judge가 pairwise verdict prompt를 안정적으로 수행하지 못한 사례로 읽는 게 맞습니다.

이 슬라이드에서 청중이 기억해야 할 건 하나입니다.
cross-family 결과를 보여주는 건 중요하지만,
그 자체가 곧 높은 품질을 의미하지는 않습니다.
특히 작은 open-weight judge를 pairwise judge로 쓸 때는 ‘모델 크기’와 ‘verdict format compliance’를 반드시 같이 봐야 합니다.""",
        ),
        SlideSpec(
            title="Phase 5: GPT-4o-mini는 외부 기준점으로 충분한가",
            section="Results",
            purpose="closed API judge를 외부 기준점으로 썼을 때 메인 결과가 얼마나 유지되는지 보여준다.",
            rq_tag="RQ3",
            bullets=[
                "Qwen32 ↔ GPT-4o-mini: Spearman ρ=0.964, Kendall τ=0.905",
                "Pairwise inconsistency도 33.99%로 Qwen32의 32.86%와 매우 유사",
                "즉 Qwen32 결과는 완전한 closed-source artifact는 아님",
            ],
            images=[FIG / "fig18_phase5_gpt4omini.png", FIG / "fig16_phase345_judge_summary.png"],
            image_captions=["summary chart: Phase 5 GPT-4o-mini 결과", "comparison chart: Phase 3/4/5 judge 요약"],
            notes="""Phase 5는 제 발표에서 외적 타당성을 가장 강하게 받쳐주는 슬라이드입니다.
GPT-4o-mini는 full GPT-4는 아니지만, 적어도 외부 API judge라는 기준점 역할은 충분히 합니다.

여기서 중요한 숫자는 Qwen32와의 rank agreement입니다.
Spearman 0.964, Kendall tau 0.905면 seen-7 모델 서열은 거의 같다고 봐도 됩니다.
그리고 pairwise inconsistency도 34% 정도라 Qwen32와 사실상 비슷한 수준입니다.

이 결과가 왜 중요하냐면,
이제 저는 Qwen32 main result를 단순히 ‘Qwen family 안에서만 그럴 수 있다’고만 말하지 않아도 되기 때문입니다.
적어도 외부 judge 하나를 붙였을 때 큰 서열 패턴은 유지된다는 걸 보였습니다.

다만 여기도 해석을 분리해야 합니다.
rank pattern consistency와 question-level decision consistency는 다릅니다.
exact pairwise agreement는 0.58 수준이기 때문에,
두 judge가 각 질문마다 똑같이 판단한다고 말하면 과장입니다.
이 슬라이드의 메시지는 ‘외부 기준점에서도 큰 흐름은 유지된다’입니다.""",
        ),
        SlideSpec(
            title="세 judge를 한 번에 비교하면 무엇이 남는가",
            section="Results",
            purpose="Phase 3/4/5를 한 장으로 묶어, 어떤 judge를 언제 쓰는 게 맞는지 결론을 준다.",
            rq_tag="RQ3",
            bullets=[
                "Qwen32와 GPT-4o-mini는 rank agreement가 매우 높고, InternLM20B는 그보다 한 단계 낮은 보조 judge다.",
                "InternLM7B는 single-grade sanity check는 가능하지만 pairwise judge로는 부적합하다.",
                "실전 권고는 `연구용 메인 judge는 32B급`, `보조 sanity check는 InternLM20B / GPT-4o-mini`다.",
            ],
            images=[FIG / "fig16_phase345_judge_summary.png"],
            image_captions=["comparison chart: Phase 3/4/5 judge 통합 요약"],
            stat_boxes=[("Qwen32↔GPT", "ρ=0.964"), ("Qwen32↔InternLM20B", "ρ=0.893"), ("GPT pairwise", "33.99%"), ("InternLM7B error", "72.6%")],
            notes="""이 슬라이드는 여러 judge 결과를 실제 의사결정 문장으로 바꾸는 슬라이드입니다.
개별 phase를 따로 보면 정보가 많은 대신, 청중은 결국 ‘그래서 무엇을 쓰면 되는데?’를 묻습니다.

제 답은 비교적 분명합니다.
메인 연구 결과를 만들 때는 32B급 judge가 필요합니다.
그리고 그 결과가 완전히 family-specific artifact가 아닌지 보려면
InternLM20B나 GPT-4o-mini 같은 보조 judge를 붙이는 게 좋습니다.

반대로 InternLM7B는 이 발표에서 일부러 숨기지 않은 실패 사례입니다.
single-grade 서열 확인은 가능하지만, pairwise verdict format 자체가 안정적이지 않아서
메인 judge로 쓰면 안 됩니다.

즉 이 슬라이드는 모델 비교 결과이면서 동시에 judge selection guideline입니다.
랩미팅에서는 이 practical takeaway가 꽤 중요하다고 생각합니다.""",
        ),
        SlideSpec(
            title="tinyMT-Bench same-set 결과: 왜 40문항이 강했는가",
            section="Results",
            purpose="문항 축소 아이디어의 출발점과 same-set 상한을 명확히 설명한다.",
            rq_tag="RQ4",
            bullets=[
                "same-set 기준 Top-Disc-40은 ρ=1.000으로 full-80 ranking을 완전 보존",
                "Top-Disc-25도 ρ=0.964로 강하지만, 1개 순위 역전이 발생",
                "same-set에서는 40문항이 가장 설득력 있는 비용 절감 upper bound",
            ],
            images=[FIG / "fig9_tiny_mt_bench.png"],
            image_captions=["line plot: same-set tinyMT-Bench 결과"],
            notes="""이 슬라이드는 문항 축소 결과의 출발점입니다.
same-set이라는 건, 문항을 고를 때 쓴 모델 집합과 검증할 때 쓴 모델 집합이 같은 경우입니다.
이 설정에서는 Top-Disc-40이 full-80 ranking을 완전히 보존했습니다.

그래서 초기에는 아주 매력적인 결과처럼 보입니다.
80문항을 40문항으로 줄이면 비용을 절반으로 줄이면서도 순위는 그대로니까요.
Top-Disc-25도 꽤 좋지만 0.964라는 건 실제로는 1개 순위 역전에 해당합니다.

다만 제가 same-set upper bound라는 표현을 계속 쓰는 이유는,
이 결과가 선택과 검증을 같은 모델 집합에서 했기 때문입니다.
즉 좋은 결과 자체는 사실이지만,
이걸 곧바로 외부 일반화로 읽으면 안 됩니다.

그래서 다음 슬라이드가 중요합니다.
same-set에서는 40문항이 예쁘게 나오지만,
hold-out으로 가면 메시지가 조금 달라집니다.
그리고 그 차이를 정직하게 보여주는 게 이 발표의 중요한 미덕이라고 생각합니다.""",
        ),
        SlideSpec(
            title="330-split hold-out은 same-set 상한을 얼마나 바꾸는가",
            section="Results",
            purpose="Phase 6가 lucky draw가 아니라 반복 hold-out 결과라는 점을 데이터로 보여준다.",
            rq_tag="RQ4",
            bullets=[
                "11개 모델에서 leave-4-out hold-out 330 split, 각 N마다 random 200개 비교",
                "Top-Disc-40 mean ρ: Qwen 0.968 [0.80, 1.00] / InternLM20B 0.922 [0.40, 1.00] / GPT-4o-mini 0.959 [0.80, 1.00]",
                "Top-Disc-60 mean ρ: 0.998 [1.00, 1.00] / 0.995 [1.00, 1.00] / 0.972 [0.80, 1.00]",
                "beats-random split rate는 40문항에서 80.6% / 78.5% / 85.5%, 60문항에서 75.8% / 95.8% / 81.2%",
                "same-set의 아름다운 상한은 40문항이지만, cross-judge 운영점은 60문항이 더 안전함",
            ],
            images=[FIG / "fig15_tiny_mt_bench_generalization.png"],
            image_captions=["line plot: repeated hold-out generalization 결과"],
            notes="""이 슬라이드가 이번 작업에서 가장 크게 업그레이드된 부분입니다.
처음엔 single split hold-out만 있었는데, 그걸로는 운 좋은 한 번이라는 비판을 피하기 어려웠습니다.
그래서 현재는 11개 모델 풀에서 4개를 test로 고르는 모든 조합, 즉 330개 split을 전부 돌렸습니다.

각 split마다 dev 7로 문항 변별도를 다시 계산해서 Top-Disc를 뽑고,
같은 N에 대해 random subset 200개도 같이 비교했습니다.
즉 이건 더 이상 pilot이 아니라 distributional result입니다.

결과를 보면 40문항도 충분히 강합니다.
평균 rho가 세 judge 모두 random 평균보다 높고, split win rate도 78~86% 정도 나옵니다.
하지만 발표에서 제가 더 강조하려는 건 60문항입니다.
40문항은 same-set의 아름다운 upper bound이고,
repeated hold-out까지 합치면 60문항이 더 안전한 cross-judge operating point로 보입니다.

이 슬라이드 덕분에 이제 tinyMT-Bench를 단순 아이디어 수준이 아니라,
상당히 방어 가능한 경험적 권고 수준까지 올릴 수 있게 됐습니다.""",
        ),
        SlideSpec(
            title="그래서 이 연구에서 무엇을 믿고, 무엇은 아직 조심해야 하나",
            section="Wrap-up",
            purpose="주장의 범위와 한계를 스스로 통제해 발표 전체 신뢰도를 끌어올린다.",
            bullets=[
                "믿을 수 있는 것: Qwen same-family scaling trend, cross-family broad rank stability, repeated hold-out에서의 60문항 권고",
                "조심할 것: 11개 모델 풀은 7B–11B same-era chat 모델 중심",
                "운영 기준선: 60문항에서는 세 judge 모두 mean ρ≥0.97, 40문항은 judge에 따라 mean ρ≈0.92까지 내려감",
                "아직 열린 질문: 70B급, code-specialist, multilingual 환경까지도 같은가",
            ],
            side_quote="좋은 발표는 강한 결과를 크게 말하고, 약한 결과를 정확히 말합니다.",
            notes="""이 슬라이드는 사실 발표 전체의 tone을 결정합니다.
결과가 아무리 좋아도 범위를 잘못 말하면 점수를 잃습니다.

제가 지금 자신 있게 말할 수 있는 건 세 가지입니다.
Qwen same-family에서는 judge가 커질수록 inconsistency가 줄었습니다.
InternLM20B와 GPT-4o-mini를 붙여도 큰 서열은 유지됩니다.
그리고 repeated hold-out 기준으로 보면 문항 축소는 60문항 정도가 가장 안전합니다.

반대로 아직 조심해야 하는 것도 분명합니다.
현재 model pool은 같은 세대의 7B–11B open-weight chat 모델에 치우쳐 있습니다.
그래서 이 결과를 70B급 모델이나 code-specialist 모델, 혹은 multilingual setting까지
바로 일반화해서 말하는 건 과장입니다.

저는 이 슬라이드가 오히려 발표의 강점이라고 생각합니다.
왜냐하면 청중은 ‘모든 걸 다 안다’고 말하는 발표보다,
‘여기까지는 강하고, 여기서부터는 아직 열린 문제다’라고 말하는 발표를 더 신뢰하기 때문입니다.""",
            stat_boxes=[("same-family", "강함"), ("cross-family", "보조 검증 완료"), ("tinyMT hold-out", "60문항 권고"), ("남은 과제", "cross-scale / multilingual")],
        ),
        SlideSpec(
            title="최종 정리: 베이스 논문을 넘어, 오픈소스 judge를 실제 도구로 읽는 법",
            section="Wrap-up",
            purpose="20분 발표의 메시지를 세 문장으로 남기고 토론으로 넘긴다.",
            bullets=[
                "베이스 논문은 GPT-4 judge의 가능성을 열었고, 내 작업은 오픈소스 judge의 실전 신뢰도를 따졌다.",
                "메인 결과는 `32B open-weight judge + abstain ensemble + 60문항 운영점`이다.",
                "Git 저장소는 논문 결과뿐 아니라 raw judgment, figure, deck까지 다시 생성 가능하게 정리했다.",
            ],
            side_quote="질문 받겠습니다.",
            stat_boxes=[("메인", "32B judge"), ("보조", "InternLM / GPT"), ("운영점", "60문항"), ("패키지", "raw→deck")],
            notes="""마지막 슬라이드에서는 세 문장만 남기겠습니다.

첫째, 베이스 논문은 GPT-4 judge가 인간 선호를 근사할 수 있다는 문을 열었습니다.
둘째, 저는 그 문을 오픈소스 judge 쪽으로 옮겨 보면서,
무엇이 실제로 믿을 만하고 무엇은 아직 조심해야 하는지 정리했습니다.
셋째, 이 저장소는 결과표만 있는 게 아니라 raw judgment부터 figure, paper, 그리고 오늘 발표자료까지 재생성 가능한 연구 자산으로 정리했습니다.

발표 이후 질문이 들어오면 보통 세 갈래일 겁니다.
왜 32B냐, 왜 60문항이냐, 그리고 왜 이 결과를 더 넓게 일반화하지 않느냐.
이 슬라이드는 그 질문들에 대한 짧은 답변이기도 합니다.

제 답은 이렇습니다.
32B는 현재 실험 범위에서 가장 좋은 open-weight judge였고,
60문항은 repeated hold-out에서 가장 안전한 운영점이었고,
그 이상의 일반화는 지금 당장 주장하지 않는 것이 더 학문적으로 정직합니다.

이상으로 발표를 마치고 질문 받겠습니다.""",
        ),
    ]


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = get_slide_specs()
    total = len(slides)
    for idx, spec in enumerate(slides, start=1):
        if idx == 1:
            build_title_slide(prs, spec, total)
        elif spec.layout == "repo":
            build_repo_slide(prs, spec, idx, total)
        elif spec.layout == "timeline":
            build_timeline_slide(prs, spec, idx, total)
        elif spec.layout == "summary":
            build_summary_slide(prs, spec, idx, total)
        else:
            build_bullets_slide(prs, spec, idx, total)

    prs.save(PPTX_PATH)
    build_notes_md(slides)
    print(PPTX_PATH)
    print(NOTES_MD_PATH)


if __name__ == "__main__":
    build_deck()
