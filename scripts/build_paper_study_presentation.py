#!/usr/bin/env python3
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
BASE_FIG = FIG / "base_paper"
OUT_DIR = ROOT / "presentation"
OUT_DIR.mkdir(exist_ok=True)
PPTX_PATH = OUT_DIR / "mt_bench_paper_study.pptx"
NOTES_MD_PATH = OUT_DIR / "mt_bench_paper_study_notes.md"


BG = RGBColor(246, 242, 235)
DEEP = RGBColor(40, 44, 58)
NAVY = RGBColor(55, 86, 103)
TEAL = RGBColor(181, 101, 78)
SEA = RGBColor(128, 151, 140)
GOLD = RGBColor(191, 154, 94)
GREEN = RGBColor(102, 129, 133)
TEXT = RGBColor(54, 58, 72)
MUTED = RGBColor(112, 118, 125)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(218, 212, 202)
PALE = RGBColor(239, 233, 225)
PALE2 = RGBColor(232, 238, 235)
PAPER_BG = RGBColor(247, 243, 237)
PAPER_RAIL = RGBColor(167, 97, 77)
PAPER_TINT = RGBColor(240, 232, 223)
RESEARCH_BG = RGBColor(239, 244, 246)
RESEARCH_RAIL = RGBColor(59, 88, 110)
RESEARCH_TINT = RGBColor(226, 234, 239)

FONT = "Apple SD Gothic Neo"
TITLE_SIZE = 26
BODY_SIZE = 17
SMALL_SIZE = 10

GRID_LEFT = 0.55
GRID_TOP = 1.42
GRID_WIDTH = 12.2
GRID_GAP = 0.28
GRID_HALF = (GRID_WIDTH - GRID_GAP) / 2
PANEL_H = 4.95


@dataclass
class SlideSpec:
    title: str
    section: str
    purpose: str
    bullets: list[str]
    notes: str
    layout: str = "content"
    takeaway: str | None = None
    images: list[Path] = field(default_factory=list)
    image_captions: list[str] = field(default_factory=list)
    stat_boxes: list[tuple[str, str]] = field(default_factory=list)
    cards: list[tuple[str, str]] = field(default_factory=list)
    compare_columns: list[tuple[str, list[str]]] = field(default_factory=list)


def section_theme(section: str) -> dict[str, RGBColor | str]:
    if section in {"Base Paper"}:
        return {
            "bg": PAPER_BG,
            "rail": PAPER_RAIL,
            "tint": PAPER_TINT,
            "footer": "논문 리뷰 | base paper and original claims",
            "section_label": "논문 리뷰",
        }
    if section in {"Reproduction", "Results", "Wrap-up", "Q&A"}:
        return {
            "bg": RESEARCH_BG,
            "rail": RESEARCH_RAIL,
            "tint": RESEARCH_TINT,
            "footer": "내 연구 | reproduction and extensions",
            "section_label": "내 연구",
        }
    return {
        "bg": BG,
        "rail": TEAL,
        "tint": PALE,
        "footer": "MT-Bench paper study | base paper + reproduction repo",
        "section_label": section,
    }


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    *,
    font_size=BODY_SIZE,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    fill=None,
    line=None,
    margin=0.07,
    valign=MSO_ANCHOR.TOP,
    auto_fit=True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line
    else:
        box.line.fill.background()
    return box


def header_title_size(text: str) -> int:
    n = len(text)
    if n >= 34:
        return 21
    if n >= 26:
        return 22
    return 24


def add_bullets(box, bullets: Iterable[str], *, font_size=BODY_SIZE, color=TEXT):
    tf = box.text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {bullet}"
        p.level = 0
        p.space_after = Pt(8)
        p.space_before = Pt(0)
        p.line_spacing = 1.12
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box


def add_picture(slide, path: Path, left, top, width=None, height=None):
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_panel(slide, left, top, width, height, *, fill=WHITE, line=LINE, margin=0.12):
    return add_textbox(slide, left, top, width, height, "", fill=fill, line=line, margin=margin)


def add_fitted_picture_panel(slide, path: Path, left, top, width, height, *, padding=0.12, fill=WHITE, line=LINE):
    add_panel(slide, left, top, width, height, fill=fill, line=line, margin=padding)
    pad = Inches(padding)
    inner_left = left + pad
    inner_top = top + pad
    inner_w = width - 2 * pad
    inner_h = height - 2 * pad

    with Image.open(path) as img:
        img_w, img_h = img.size

    scale = min(inner_w / img_w, inner_h / img_h)
    draw_w = int(img_w * scale)
    draw_h = int(img_h * scale)
    pic_left = int(left + (width - draw_w) / 2)
    pic_top = int(top + (height - draw_h) / 2)
    return add_picture(slide, path, pic_left, pic_top, width=draw_w, height=draw_h)


def add_stat_box(slide, left, top, label, value, accent):
    add_textbox(slide, left, top, Inches(2.22), Inches(0.96), "", fill=WHITE, line=LINE, margin=0.06)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left + Inches(0.12),
        top + Inches(0.12),
        Inches(0.55),
        Inches(0.2),
    )
    chip = slide.shapes[-1]
    chip.fill.solid()
    chip.fill.fore_color.rgb = accent
    chip.line.fill.background()
    add_textbox(
        slide,
        left + Inches(0.14),
        top + Inches(0.12),
        Inches(1.85),
        Inches(0.2),
        label,
        font_size=9,
        bold=True,
        color=WHITE,
        margin=0.01,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        left + Inches(0.12),
        top + Inches(0.38),
        Inches(1.9),
        Inches(0.34),
        value,
        font_size=19,
        bold=True,
        color=DEEP,
        margin=0.01,
    )


def paint_background(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_header(slide, title: str, section: str, idx: int, total: int):
    theme = section_theme(section)
    paint_background(slide, theme["bg"])
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.36), Inches(0.42), Inches(0.14), Inches(0.92))
    rail = slide.shapes[-1]
    rail.fill.solid()
    rail.fill.fore_color.rgb = theme["rail"]
    rail.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.36), Inches(0.2), Inches(12.45), Inches(0.04))
    top_line = slide.shapes[-1]
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = LINE
    top_line.line.fill.background()
    add_textbox(
        slide,
        Inches(0.68),
        Inches(0.36),
        Inches(9.4),
        Inches(0.58),
        title,
        font_size=header_title_size(title),
        bold=True,
        color=DEEP,
        margin=0.0,
    )
    add_textbox(
        slide,
        Inches(10.5),
        Inches(0.44),
        Inches(1.55),
        Inches(0.22),
        theme["section_label"],
        font_size=10,
        bold=True,
        color=DEEP,
        align=PP_ALIGN.CENTER,
        fill=theme["tint"],
        line=theme["tint"],
        margin=0.01,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        Inches(12.1),
        Inches(0.46),
        Inches(0.75),
        Inches(0.18),
        f"{idx}/{total}",
        font_size=10,
        bold=True,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        margin=0.0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(7.02), Inches(12.25), Inches(0.04))
    foot = slide.shapes[-1]
    foot.fill.solid()
    foot.fill.fore_color.rgb = LINE
    foot.line.fill.background()
    add_textbox(
        slide,
        Inches(0.55),
        Inches(7.06),
        Inches(5.1),
        Inches(0.18),
        theme["footer"],
        font_size=8,
        color=MUTED,
        margin=0.0,
    )


def add_purpose_box(slide, purpose: str, section: str):
    theme = section_theme(section)
    add_textbox(
        slide,
        Inches(0.48),
        Inches(0.76),
        Inches(12.35),
        Inches(0.4),
        f"이 슬라이드의 목적: {purpose}",
        font_size=12,
        bold=True,
        color=DEEP,
        fill=theme["tint"],
        line=theme["tint"],
        margin=0.03,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build_title_slide(prs: Presentation, spec: SlideSpec, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide, PAPER_BG)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(6.86), Inches(7.5))
    left_panel = slide.shapes[-1]
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = PAPER_BG
    left_panel.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.86), Inches(0), Inches(6.47), Inches(7.5))
    right_panel = slide.shapes[-1]
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RESEARCH_BG
    right_panel.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.28), Inches(0.55), Inches(0.16), Inches(5.45))
    accent = slide.shapes[-1]
    accent.fill.solid()
    accent.fill.fore_color.rgb = PAPER_RAIL
    accent.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.12), Inches(0.72), Inches(0.14), Inches(4.9))
    accent2 = slide.shapes[-1]
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = RESEARCH_RAIL
    accent2.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(6.28), Inches(12.08), Inches(0.82))
    bottom = slide.shapes[-1]
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = DEEP
    bottom.line.fill.background()

    add_textbox(
        slide,
        Inches(0.78),
        Inches(0.82),
        Inches(5.88),
        Inches(1.7),
        "Judging LLM-as-a-Judge\nwith MT-Bench\nand Chatbot Arena",
        font_size=27,
        bold=True,
        color=DEEP,
        margin=0.0,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.78),
        Inches(5.3),
        Inches(0.28),
        "논문 스터디",
        font_size=14,
        color=PAPER_RAIL,
        bold=True,
        margin=0.0,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.2),
        Inches(5.55),
        Inches(1.05),
        "앞 절반은 원 논문의 설계와 메시지를 읽고,\n뒤 절반은 그 프로토콜을 내 오픈소스 judge 실험으로 어떻게 옮겼는지 설명합니다.",
        font_size=16,
        color=TEXT,
    )
    add_textbox(slide, Inches(7.42), Inches(0.9), Inches(4.9), Inches(0.26), "Base paper figures that anchor the talk", font_size=14, bold=True, color=RESEARCH_RAIL, margin=0.0)
    add_picture(slide, BASE_FIG / "paper_chatbot_arena_ui.png", Inches(7.38), Inches(1.3), width=Inches(5.25))
    add_picture(slide, BASE_FIG / "paper_table8_scores.png", Inches(7.72), Inches(4.22), width=Inches(4.55))
    add_textbox(slide, Inches(7.42), Inches(4.0), Inches(5.18), Inches(0.18), "Original paper Figure 19", font_size=9, color=MUTED, align=PP_ALIGN.RIGHT, margin=0.0)
    add_textbox(slide, Inches(7.78), Inches(6.0), Inches(4.42), Inches(0.18), "Original paper Table 8", font_size=9, color=MUTED, align=PP_ALIGN.RIGHT, margin=0.0)

    add_textbox(
        slide,
        Inches(0.84),
        Inches(6.46),
        Inches(11.4),
        Inches(0.26),
        "핵심 질문: 이 논문의 strong-judge 메시지가 내 오픈소스 judge 재현 실험에서도 어디까지 유지되는가?",
        font_size=15,
        bold=True,
        color=WHITE,
        margin=0.0,
    )
    add_textbox(
        slide,
        Inches(0.84),
        Inches(6.82),
        Inches(6.4),
        Inches(0.18),
        "박승현 | CLINK Lab | paper study + my reproduction",
        font_size=10,
        color=PALE2,
        margin=0.0,
    )
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_divider_slide(prs: Presentation, spec: SlideSpec, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = section_theme(spec.section)
    paint_background(slide, theme["bg"])
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(12.2), Inches(0.04))
    top = slide.shapes[-1]
    top.fill.solid()
    top.fill.fore_color.rgb = LINE
    top.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.08), Inches(0.14), Inches(3.15))
    stripe = slide.shapes[-1]
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = theme["rail"]
    stripe.line.fill.background()
    number = spec.title.split(".")[0] if "." in spec.title else str(idx)
    add_textbox(
        slide,
        Inches(9.35),
        Inches(1.15),
        Inches(2.9),
        Inches(1.55),
        number,
        font_size=74,
        bold=True,
        color=LINE,
        align=PP_ALIGN.RIGHT,
        margin=0.0,
    )

    add_textbox(
        slide,
        Inches(0.94),
        Inches(1.1),
        Inches(6.8),
        Inches(0.38),
        theme["section_label"],
        font_size=15,
        bold=True,
        color=theme["rail"],
        margin=0.0,
    )
    add_textbox(
        slide,
        Inches(0.94),
        Inches(1.7),
        Inches(7.6),
        Inches(0.95),
        spec.title,
        font_size=28,
        bold=True,
        color=DEEP,
        margin=0.0,
    )
    add_textbox(
        slide,
        Inches(0.94),
        Inches(3.02),
        Inches(7.45),
        Inches(1.05),
        spec.takeaway or "",
        font_size=18,
        color=DEEP,
        fill=theme["tint"],
        line=theme["tint"],
        margin=0.12,
    )
    add_textbox(
        slide,
        Inches(12.0),
        Inches(0.48),
        Inches(0.72),
        Inches(0.16),
        f"{idx}/{total}",
        font_size=10,
        bold=True,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        margin=0.0,
    )
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_overview_slide(prs: Presentation, spec: SlideSpec, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = section_theme(spec.section)
    add_header(slide, spec.title, spec.section, idx, total)
    add_purpose_box(slide, spec.purpose, spec.section)
    left = add_textbox(slide, Inches(0.55), Inches(1.42), Inches(7.0), Inches(4.95), "", fill=WHITE, line=LINE, margin=0.12)
    add_bullets(left, spec.bullets, font_size=20)
    add_textbox(
        slide,
        Inches(7.85),
        Inches(1.42),
        Inches(4.95),
        Inches(0.32),
        "오늘 발표를 듣는 포인트",
        font_size=18,
        bold=True,
        color=DEEP,
    )
    insight = add_textbox(slide, Inches(7.85), Inches(1.86), Inches(4.95), Inches(3.9), "", fill=theme["tint"], line=theme["tint"], margin=0.14)
    add_bullets(
        insight,
        [
            "원 논문이 왜 영향력이 컸는지 먼저 이해한다.",
            "내 재현은 어디까지가 faithful reproduction이고 어디부터가 extension인지 구분한다.",
            "결과는 same-family / cross-family / hold-out 순서로 보수적으로 읽는다.",
        ],
        font_size=17,
    )
    if spec.takeaway:
        add_textbox(
            slide,
            Inches(7.85),
            Inches(6.0),
            Inches(4.95),
            Inches(0.65),
            spec.takeaway,
            font_size=16,
            bold=True,
            color=DEEP,
            fill=theme["tint"],
            line=theme["tint"],
            valign=MSO_ANCHOR.MIDDLE,
        )
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_cards_slide(prs: Presentation, spec: SlideSpec, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = section_theme(spec.section)
    add_header(slide, spec.title, spec.section, idx, total)
    add_purpose_box(slide, spec.purpose, spec.section)

    if spec.takeaway:
        add_textbox(
            slide,
            Inches(0.55),
            Inches(1.42),
            Inches(12.2),
            Inches(0.62),
            spec.takeaway,
            font_size=19,
            bold=True,
            color=DEEP,
            fill=theme["tint"],
            line=theme["tint"],
            valign=MSO_ANCHOR.MIDDLE,
        )
        card_top = Inches(2.24)
    else:
        card_top = Inches(1.42)

    count = max(1, len(spec.cards))
    if count == 4:
        card_h = 1.58
        row_gap = 0.2
        row_tops = [card_top, card_top + Inches(card_h + row_gap)]
        col_lefts = [Inches(GRID_LEFT), Inches(GRID_LEFT + GRID_HALF + GRID_GAP)]
        card_w = Inches(GRID_HALF)
        for i, (head, body) in enumerate(spec.cards):
            row = i // 2
            col = i % 2
            left = col_lefts[col]
            top = row_tops[row]
            add_panel(slide, left, top, card_w, Inches(card_h), fill=WHITE, line=LINE, margin=0.12)
            add_textbox(slide, left + Inches(0.1), top + Inches(0.12), card_w - Inches(0.2), Inches(0.22), head, font_size=15, bold=True, color=theme["rail"], margin=0.0)
            add_textbox(slide, left + Inches(0.1), top + Inches(0.44), card_w - Inches(0.2), Inches(0.88), body, font_size=13.5, color=TEXT, margin=0.0)
        bullet_top = row_tops[-1] + Inches(card_h + 0.18)
        bullet_h = Inches(0.82)
    else:
        gap = 0.22
        width = (GRID_WIDTH - gap * (count - 1)) / count
        card_w = Inches(width)
        for i, (head, body) in enumerate(spec.cards):
            left = Inches(GRID_LEFT + i * (width + gap))
            add_panel(slide, left, card_top, card_w, Inches(2.22), fill=WHITE, line=LINE, margin=0.12)
            add_textbox(slide, left + Inches(0.08), card_top + Inches(0.12), card_w - Inches(0.16), Inches(0.26), head, font_size=15, bold=True, color=theme["rail"], margin=0.0)
            add_textbox(slide, left + Inches(0.08), card_top + Inches(0.52), card_w - Inches(0.16), Inches(1.3), body, font_size=14, color=TEXT, margin=0.0)
        bullet_top = Inches(4.82)
        bullet_h = Inches(1.42)

    if spec.bullets:
        body = add_textbox(slide, Inches(GRID_LEFT), bullet_top, Inches(GRID_WIDTH), bullet_h, "", fill=WHITE, line=LINE, margin=0.13)
        add_bullets(body, spec.bullets, font_size=15 if count == 4 else 16)
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_compare_slide(prs: Presentation, spec: SlideSpec, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = section_theme(spec.section)
    add_header(slide, spec.title, spec.section, idx, total)
    add_purpose_box(slide, spec.purpose, spec.section)

    if spec.takeaway:
        add_textbox(
            slide,
            Inches(0.55),
            Inches(1.42),
            Inches(12.2),
            Inches(0.54),
            spec.takeaway,
            font_size=17,
            bold=True,
            color=DEEP,
            fill=theme["tint"],
            line=theme["tint"],
            valign=MSO_ANCHOR.MIDDLE,
        )
        top = Inches(2.16)
    else:
        top = Inches(1.42)

    for i, (head, items) in enumerate(spec.compare_columns[:2]):
        left = Inches(GRID_LEFT + i * (GRID_HALF + GRID_GAP))
        add_panel(slide, left, top, Inches(GRID_HALF), Inches(4.78), fill=WHITE, line=LINE, margin=0.13)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.12), Inches(GRID_HALF - 0.2), Inches(0.26), head, font_size=18, bold=True, color=DEEP, margin=0.0)
        body = add_textbox(slide, left + Inches(0.1), top + Inches(0.52), Inches(GRID_HALF - 0.2), Inches(3.9), "", fill=None, line=None, margin=0.0)
        add_bullets(body, items, font_size=17)
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_content_slide(prs: Presentation, spec: SlideSpec, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = section_theme(spec.section)
    add_header(slide, spec.title, spec.section, idx, total)
    add_purpose_box(slide, spec.purpose, spec.section)
    total_chars = sum(len(b) for b in spec.bullets)

    if len(spec.images) == 2:
        panel_top = Inches(1.38)
        panel_h = Inches(3.18)
        left_x = Inches(GRID_LEFT)
        right_x = Inches(GRID_LEFT + GRID_HALF + GRID_GAP)
        add_fitted_picture_panel(slide, spec.images[0], left_x, panel_top, Inches(GRID_HALF), panel_h)
        add_fitted_picture_panel(slide, spec.images[1], right_x, panel_top, Inches(GRID_HALF), panel_h)
        if spec.image_captions:
            add_textbox(slide, left_x, Inches(4.68), Inches(GRID_HALF), Inches(0.22), spec.image_captions[0], font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
            add_textbox(slide, right_x, Inches(4.68), Inches(GRID_HALF), Inches(0.22), spec.image_captions[1], font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
        body = add_textbox(slide, Inches(GRID_LEFT), Inches(4.96), Inches(GRID_WIDTH), Inches(1.48), "", fill=WHITE, line=LINE, margin=0.13)
        add_bullets(body, spec.bullets, font_size=15 if total_chars > 260 else 16)
        if spec.stat_boxes:
            for i, (label, value) in enumerate(spec.stat_boxes):
                accent = [NAVY, TEAL, GOLD, GREEN][i % 4]
                add_stat_box(slide, Inches(0.62 + i * 2.42), Inches(6.1), label, value, accent)
        elif spec.takeaway:
            add_textbox(
                slide,
                Inches(GRID_LEFT),
                Inches(6.12),
                Inches(GRID_WIDTH),
                Inches(0.6),
                spec.takeaway,
                font_size=16,
                bold=True,
                color=DEEP,
                fill=theme["tint"],
                line=theme["tint"],
                valign=MSO_ANCHOR.MIDDLE,
            )
    elif len(spec.images) == 1:
        image_left = Inches(0.45)
        image_top = Inches(1.38)
        image_w = Inches(7.72)
        image_h = Inches(4.35)
        add_fitted_picture_panel(slide, spec.images[0], image_left, image_top, image_w, image_h)
        if spec.image_captions:
            add_textbox(slide, Inches(0.62), Inches(5.88), Inches(7.48), Inches(0.22), spec.image_captions[0], font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
        body = add_textbox(slide, Inches(8.38), Inches(1.38), Inches(4.35), Inches(4.35), "", fill=WHITE, line=LINE, margin=0.13)
        add_bullets(body, spec.bullets, font_size=15 if total_chars > 240 else 16)
        if spec.stat_boxes:
            for i, (label, value) in enumerate(spec.stat_boxes):
                accent = [NAVY, TEAL, GOLD, GREEN][i % 4]
                x = 8.45 + (i % 2) * 2.12
                y = 5.94 + (i // 2) * 0.98
                add_stat_box(slide, Inches(x), Inches(y), label, value, accent)
        elif spec.takeaway:
            add_textbox(
                slide,
                Inches(8.38),
                Inches(5.96),
                Inches(4.35),
                Inches(0.75),
                spec.takeaway,
                font_size=15,
                bold=True,
                color=DEEP,
                fill=theme["tint"],
                line=theme["tint"],
                valign=MSO_ANCHOR.MIDDLE,
            )
    else:
        body = add_textbox(slide, Inches(0.55), Inches(1.42), Inches(7.65), Inches(4.95), "", fill=WHITE, line=LINE, margin=0.14)
        add_bullets(body, spec.bullets, font_size=15 if total_chars > 320 else 16)
        add_textbox(slide, Inches(8.45), Inches(1.42), Inches(4.3), Inches(0.32), "핵심 해석", font_size=18, bold=True, color=DEEP)
        right = add_textbox(slide, Inches(8.45), Inches(1.8), Inches(4.3), Inches(4.57), "", fill=theme["tint"], line=theme["tint"], margin=0.14)
        if spec.takeaway:
            add_textbox(slide, Inches(8.62), Inches(2.02), Inches(3.9), Inches(1.28), spec.takeaway, font_size=16, bold=True, color=DEEP, margin=0.0)
        if spec.stat_boxes:
            for i, (label, value) in enumerate(spec.stat_boxes):
                accent = [NAVY, TEAL, GOLD, GREEN][i % 4]
                x = 8.62 + (i % 2) * 2.02
                y = 3.62 + (i // 2) * 1.0
                add_stat_box(slide, Inches(x), Inches(y), label, value, accent)
        else:
            add_textbox(slide, Inches(8.62), Inches(3.82), Inches(3.78), Inches(1.6), "이 슬라이드는 세부 실험보다도\n왜 이 논점이 다음 슬라이드의 전제가 되는지\n이해시키는 역할을 합니다.", font_size=14, color=TEXT)
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_timeline_slide(prs: Presentation, spec: SlideSpec, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = section_theme(spec.section)
    add_header(slide, spec.title, spec.section, idx, total)
    add_purpose_box(slide, spec.purpose, spec.section)

    phases = [
        ("P1", "Self-judge"),
        ("P2", "Sanity"),
        ("P3", "Scaling"),
        ("P4", "InternLM"),
        ("P5", "GPT-4o-mini"),
        ("P6", "Hold-out"),
    ]
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(2.1), Inches(11.1), Inches(0.06))
    line = slide.shapes[-1]
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    xs = [1.0, 2.95, 4.9, 6.85, 8.8, 10.75]
    accents = [GOLD, TEAL, NAVY, GREEN, DEEP, NAVY]
    for x, (tag, label), accent in zip(xs, phases, accents):
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(1.72), Inches(0.72), Inches(0.72))
        circ = slide.shapes[-1]
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent
        circ.line.fill.background()
        add_textbox(slide, Inches(x), Inches(1.91), Inches(0.72), Inches(0.15), tag, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0.0)
        add_textbox(slide, Inches(x - 0.18), Inches(2.48), Inches(1.08), Inches(0.35), label, font_size=12, bold=True, color=DEEP, align=PP_ALIGN.CENTER, margin=0.0)

    body = add_textbox(slide, Inches(GRID_LEFT), Inches(3.18), Inches(GRID_WIDTH), Inches(3.08), "", fill=WHITE, line=LINE, margin=0.13)
    add_bullets(body, spec.bullets, font_size=17)
    if spec.takeaway:
        add_textbox(slide, Inches(GRID_LEFT), Inches(6.38), Inches(GRID_WIDTH), Inches(0.42), spec.takeaway, font_size=15, bold=True, color=DEEP, fill=theme["tint"], line=theme["tint"], valign=MSO_ANCHOR.MIDDLE)
    slide.notes_slide.notes_text_frame.text = spec.notes


def build_notes_md(slides: list[SlideSpec]):
    lines = ["# MT-Bench Paper Study Notes", ""]
    for idx, spec in enumerate(slides, 1):
        lines.append(f"## Slide {idx}. {spec.title}")
        lines.append(f"- Section: {spec.section}")
        lines.append(f"- Purpose: {spec.purpose}")
        lines.append(f"- Layout: {spec.layout}")
        if spec.images:
            lines.append("- Visuals:")
            for img, cap in zip(spec.images, spec.image_captions or [p.name for p in spec.images]):
                lines.append(f"  - {img.name}: {cap}")
        if spec.bullets:
            lines.append("- Slide bullets:")
            for bullet in spec.bullets:
                lines.append(f"  - {bullet}")
        if spec.takeaway:
            lines.append(f"- Takeaway: {spec.takeaway}")
        if spec.stat_boxes:
            lines.append("- Stat boxes:")
            for label, value in spec.stat_boxes:
                lines.append(f"  - {label}: {value}")
        lines.append("")
        lines.append("### Speaker Notes")
        lines.append(spec.notes.strip())
        lines.append("")
    NOTES_MD_PATH.write_text("\n".join(lines), encoding="utf-8")


def get_slide_specs() -> list[SlideSpec]:
    return [
        SlideSpec(
            title="MT-Bench / Chatbot Arena\n논문 리뷰와 내 연구",
            section="Title",
            purpose="발표의 범위를 논문 스터디 중심으로 잡고, 뒤에 재현 결과가 이어진다는 점을 알려준다.",
            bullets=[],
            layout="title",
            notes="""오늘 발표는 랩미팅 발표보다는 논문 스터디에 더 가깝게 구성했습니다.
앞 절반은 Zheng et al.의 MT-Bench / Chatbot Arena 논문을 읽는 시간이 되고,
뒤 절반은 제가 그 프로토콜을 오픈소스 judge 파이프라인으로 어떻게 옮겼는지 설명하는 구조입니다.

따라서 이 발표를 들을 때 중요한 기준은 두 가지입니다.
첫째, 원 논문이 실제로 무엇을 주장했는지 정확히 이해하는 것.
둘째, 제 재현 실험이 그 주장 중 어떤 부분을 지지하고 어떤 부분을 보수적으로 다시 해석하는지 보는 것입니다.

오늘의 최종 메시지는 간단합니다.
MT-Bench는 여전히 좋은 benchmark이고, LLM-as-a-Judge는 유용한 접근입니다.
다만 오픈소스 judge에서는 크기, 아키텍처, pairwise format 안정성, hold-out 일반화까지 함께 봐야 합니다.""",
        ),
        SlideSpec(
            title="오늘 발표는 1부 논문 리뷰, 2부 내 연구입니다",
            section="Roadmap",
            purpose="20분 동안 어떤 질문을 어떤 순서로 볼지 먼저 정렬한다.",
            bullets=[],
            layout="cards",
            takeaway="발표를 관통하는 질문: 원 논문의 85% judge-human agreement라는 메시지는 오픈소스 judge 환경에서도 어느 정도 유지되는가?",
            cards=[
                ("1부. 논문 리뷰", "MT-Bench와 Chatbot Arena가 왜 같이 설계됐는지, 그리고 GPT-4 judge가 무엇을 설득했는지 먼저 읽습니다."),
                ("2부. 내 연구", "같은 judge protocol을 제 저장소에서 Qwen · InternLM · GPT-4o-mini 조합으로 옮긴 방식을 설명합니다."),
                ("발표의 최종 메시지", "judge scaling, cross-family sanity, repeated hold-out을 순서대로 보며 어디까지 믿을 수 있는지 정리합니다."),
            ],
            notes="""이 슬라이드는 오늘 발표의 독해법을 정하는 슬라이드입니다.
먼저 원 논문을 이해해야 뒤의 재현 실험이 단순 숫자 나열이 아니라는 점이 보입니다.

발표는 목차를 나열하기보다 하나의 질문을 세 번 나눠 답하는 구조로 듣는 게 좋습니다.
첫째, 원 논문의 주장이 정확히 무엇이었는가.
둘째, 그 프로토콜을 오픈소스 judge 환경으로 옮기면 무엇이 달라지는가.
셋째, 유지되는 결론과 더 이상 강하게 말하면 안 되는 결론은 무엇인가.

발표가 끝나면 세 질문에 답할 수 있으면 됩니다.
원 논문이 왜 중요했는가.
내 저장소가 무엇을 재현 가능하게 만들었는가.
그리고 오픈소스 judge를 실제로 어디까지 믿어도 되는가.""",
        ),
        SlideSpec(
            title="1-1. 왜 MT-Bench가 필요했는가",
            section="Base Paper",
            purpose="원 논문이 기존 벤치마크의 어떤 한계를 문제 삼았는지 잡는다.",
            bullets=[
                "객관식/단답형 중심 벤치마크는 open-ended, multi-turn 대화 품질을 충분히 반영하지 못한다.",
                "사람 선호를 직접 모으는 평가는 비싸고 느려서 반복적인 모델 개발 루프에 맞지 않는다.",
                "그래서 원 논문은 ‘사람 선호를 근사하는 자동 judge’를 중심 문제로 세운다.",
            ],
            images=[BASE_FIG / "paper_chatbot_arena_ui.png", BASE_FIG / "paper_mtbench_winrate_fig3.png"],
            image_captions=["Original paper Figure 19: Chatbot Arena UI", "Original paper Figure 3: MT-Bench 평균 승률 곡선"],
            takeaway="원 논문의 문제의식은 ‘새 benchmark 하나 만들기’보다 ‘사람 평가를 대체할 실용적 judge 만들기’에 더 가깝습니다.",
            notes="""이 논문의 출발점은 매우 현실적입니다.
대화형 모델을 평가하려면 open-ended 품질을 봐야 하는데, 기존 객관식 벤치마크는 그걸 충분히 담지 못합니다.
그렇다고 사람 평가를 계속 붙이면 속도가 너무 느리고 비용이 큽니다.

그래서 원 논문은 두 축을 동시에 만듭니다.
하나는 controlled benchmark인 MT-Bench이고,
다른 하나는 실제 사용자의 선호 데이터를 모으는 Chatbot Arena입니다.

왼쪽 그림은 Arena의 실제 인터페이스입니다.
사용자는 모델 이름을 모른 채 두 응답 중 더 나은 답변을 고릅니다.
오른쪽 그림은 MT-Bench 점수와 Arena 승률의 관계를 보여줍니다.
즉 원 논문은 benchmark와 in-the-wild preference를 따로따로 두지 않고,
둘의 연결을 통해 judge의 타당성을 설득하려 했습니다.""",
        ),
        SlideSpec(
            title="1-2. MT-Bench와 Chatbot Arena는 어떻게 역할을 나눴는가",
            section="Base Paper",
            purpose="MT-Bench와 Arena가 각각 무엇을 측정하고 어떻게 보완하는지 구조적으로 설명한다.",
            bullets=[
                "MT-Bench는 8개 카테고리, 80문항의 multi-turn benchmark로 controlled comparison을 담당한다.",
                "Chatbot Arena는 실제 사용자 pairwise preference를 통해 ecological validity를 제공한다.",
                "즉 MT-Bench는 재현 가능성, Arena는 현실 적합성을 제공하는 쌍으로 설계되었다.",
            ],
            stat_boxes=[("MT-Bench", "80문항"), ("Category", "8개"), ("Arena", "익명 pairwise"), ("Goal", "human preference")],
            takeaway="논문의 힘은 MT-Bench 하나가 아니라, benchmark와 crowd preference를 함께 설계했다는 데 있습니다.",
            notes="""MT-Bench와 Arena는 겉으로 보기엔 둘 다 평가 데이터처럼 보이지만 역할이 다릅니다.
MT-Bench는 controlled benchmark입니다. 같은 질문 세트에 대해 여러 모델을 반복적으로 비교할 수 있고,
category별 분석도 가능합니다.

반면 Arena는 통제된 benchmark가 아니라 실제 사용자 환경입니다.
질문도 자유롭고, 사용자는 두 답변을 블라인드로 비교합니다.
이건 benchmark의 재현성은 떨어지지만 현실의 인간 선호를 훨씬 더 잘 반영합니다.

원 논문이 강한 이유는 이 둘을 함께 사용했다는 점입니다.
즉 MT-Bench 점수가 높으면 Arena 승률도 높다는 관계를 보이면서,
benchmark 점수가 단순 인공적인 숫자가 아니라는 걸 설득합니다.

이 발표 뒤쪽에서 제 재현 실험은 MT-Bench 쪽을 중심으로 따라갑니다.
Arena 자체를 재현한 건 아니고, MT-Bench judge reliability를 더 자세히 캐보는 방향으로 확장했다고 보시면 됩니다.""",
        ),
        SlideSpec(
            title="1-3. 원 논문은 judge를 어떻게 썼는가",
            section="Base Paper",
            purpose="single, pairwise, reference-guided라는 judge 프로토콜을 분명히 해 둔다.",
            bullets=[
                "논문은 이 세 judge mode를 동시에 운용하면서, 위치·장문·self-enhancement bias를 별도 분석 대상으로 둔다.",
            ],
            layout="cards",
            cards=[
                ("Single-answer grading", "답변 하나를 0–10 점수로 채점합니다. 서열 안정성과 카테고리별 score range를 보기 좋습니다."),
                ("Pairwise comparison", "A/B 두 응답을 비교해 더 나은 답변을 고릅니다. 실제 사용자의 pairwise preference와 가장 닮은 setting입니다."),
                ("Reference-guided grading", "수학·코딩처럼 정답 기준이 있는 경우 reference answer를 함께 주어 채점 근거를 강화합니다."),
            ],
            takeaway="원 논문의 결론은 ‘GPT-4가 완벽하다’가 아니라, ‘강한 judge가 인간 선호를 실용적으로 근사한다’입니다.",
            notes="""이 슬라이드는 발표 전체에서 아주 중요합니다.
왜냐하면 뒤의 제 실험도 결국 같은 judge 프로토콜을 구현하고 있기 때문입니다.

Single-answer grading은 0에서 10 사이 점수를 주는 방식이고,
pairwise는 두 응답 중 winner를 직접 고르는 방식입니다.
Reference-guided grading은 특히 수학과 코딩처럼 정답 기준이 필요한 문항에서 중요합니다.

많은 사람이 원 논문을 읽으면서 GPT-4 judge가 그냥 채점기처럼 등장했다고 생각하는데,
실제로는 prompt protocol이 꽤 정교합니다.
그리고 저자들도 bias 문제를 적극적으로 분석합니다.
verbosity bias, position bias, self-enhancement bias를 같이 본다는 점은
원 논문이 생각보다 훨씬 신중한 paper라는 뜻입니다.

이 점이 나중에 제 실험에서도 중요한 기준이 됩니다.
즉 judge 성능을 볼 때 agreement만 보는 게 아니라, 남는 bias의 구조까지 같이 봐야 한다는 겁니다.""",
        ),
        SlideSpec(
            title="1-4. judge 점수평가는 실제로 어떻게 흘러가는가",
            section="Base Paper",
            purpose="답변 생성부터 parse·aggregate까지 judge scoring flow를 한 번에 보여준다.",
            bullets=[
                "generation 단계와 judge 단계의 temperature를 분리해, 답변 다양성과 평가 결정성을 서로 다른 층으로 관리한다.",
                "single-grade는 turn1 독립 prompt와 turn2 multi-turn prompt를 따로 써서 turn2 맥락 손실을 막는다.",
                "pairwise의 핵심은 winner를 한 번 맞히는 게 아니라, AB/BA swap 뒤에도 유지되는 winner만 채택하는 것이다.",
            ],
            layout="cards",
            cards=[
                ("Step 1. 답변 생성", "각 모델이 80문항, 2-turn 답변을 생성합니다. generation은 temperature 0.7로 두어 응답 다양성을 확보합니다."),
                ("Step 2. judge 호출", "single은 turn1/turn2를 채점하고, pairwise는 AB와 BA 두 순서로 각각 judge를 호출합니다. judge는 temperature 0.0입니다."),
                ("Step 3. 파싱", "single은 1–10 점수를 추출하고 실패 시 -1.0으로 기록합니다. pairwise는 A/B/tie/error를 추출합니다."),
                ("Step 4. 집계", "single은 문항별 turn 평균을 모델 점수로, pairwise는 swap 뒤에도 일치할 때만 winner로 채택합니다."),
            ],
            takeaway="이 저장소의 재현성 핵심은 prompt를 흉내 낸 것보다도, parse rule과 aggregation rule을 논문과 맞춘 데 있습니다.",
            notes="""이 슬라이드는 원 논문과 제 저장소 사이를 연결하는 가장 중요한 프로토콜 슬라이드입니다.
청중이 여기서 judge가 실제로 무엇을 입력받고 무엇을 출력하는지 이해해야,
뒤의 scaling과 abstain 결과도 설득력 있게 들립니다.

첫 단계는 답변 생성입니다. 각 모델이 80문항 2-turn 답변을 만들고,
여기에는 temperature 0.7을 둬서 응답 다양성을 허용합니다.
두 번째 단계는 judge 호출입니다. judge는 반대로 temperature 0.0으로 고정해서
평가 noise를 줄입니다.

single-grade에서는 turn1과 turn2를 같은 방식으로 다루지 않습니다.
turn2는 q1, a1, q2, a2를 함께 넣어 multi-turn context로 채점해야
실제 MT-Bench 프로토콜과 맞습니다.

pairwise에서는 AB와 BA를 둘 다 돌립니다.
핵심은 두 번의 호출을 했다는 사실이 아니라, 그 둘을 합칠 때
swap 뒤에도 같은 winner가 유지될 때만 채택한다는 점입니다.
즉 이 연구는 단순 prompt demo가 아니라, parse와 aggregate rule까지 포함한 protocol reproduction입니다.""",
        ),
        SlideSpec(
            title="1-5. 원 논문의 두 핵심 주장은 무엇이었나",
            section="Base Paper",
            purpose="뒤 재현 실험의 기준점이 되는 원 논문의 핵심 claim 두 개를 먼저 또렷하게 세운다.",
            bullets=[
                "핵심 주장 1: GPT-4 judge는 MT-Bench에서 인간 expert와 non-tie agreement 85%, human-human 81% 수준으로 맞는다.",
                "핵심 주장 2: MT-Bench 점수는 GPT-4 8.99, GPT-3.5 7.94, Vicuna-13B 6.39, LLaMA-13B 2.61처럼 모델 품질 차이를 분명하게 서열화한다.",
                "즉 원 논문은 ‘judge가 인간과 비슷하다’와 ‘점수가 모델 품질을 반영한다’를 동시에 주장한다.",
            ],
            images=[BASE_FIG / "paper_mtbench_agreement_table5.png", BASE_FIG / "paper_table8_scores.png"],
            image_captions=["Original paper Table 5: judge-human agreement", "Original paper Table 8: MT-Bench model scores"],
            takeaway="뒤의 재현 실험은 결국 이 두 문장을 다시 묻는 작업입니다. agreement가 유지되는가, 그리고 서열이 설득력 있는가?",
            notes="""이 슬라이드는 원 논문을 조금 더 공정하게 읽기 위해 넣었습니다.
원 논문을 처음 읽는 사람에게는 이 슬라이드가 가장 중요합니다.
이 논문이 실제로 주장한 핵심은 두 문장으로 요약할 수 있습니다.

첫째, GPT-4 judge는 인간 expert와 꽤 잘 맞는다.
둘째, MT-Bench 점수는 모델 품질 차이를 꽤 설득력 있게 서열화한다.

왼쪽 Table 5가 첫 번째 주장이고, 오른쪽 Table 8이 두 번째 주장입니다.
이 두 숫자가 붙으면서 LLM-as-a-Judge라는 메시지가 힘을 얻습니다.

뒤쪽 제 재현 실험도 사실은 이 두 질문을 반복합니다.
agreement가 오픈소스 judge에서도 유지되는가,
그리고 서열이 여전히 설득력 있는가.
이 슬라이드가 그 기준점입니다.""",
        ),
        SlideSpec(
            title="1-6. 원 논문은 bias를 무시하지 않았다",
            section="Base Paper",
            purpose="원 논문이 숫자만 좋은 게 아니라 judge 자체를 편향이 있는 시스템으로 다뤘다는 점을 짚는다.",
            bullets=[],
            layout="compare",
            compare_columns=[
                ("원 논문이 직접 본 편향", [
                    "position bias: 응답 순서가 바뀌면 winner가 달라질 수 있다.",
                    "verbosity bias: 더 길고 장황한 답변이 유리해질 수 있다.",
                    "self-enhancement bias: judge가 자기 패밀리 답변을 더 좋게 볼 수 있다.",
                ]),
                ("왜 내 재현과 이어지나", [
                    "agreement가 높아도 남는 오류 구조는 따로 봐야 한다.",
                    "오픈소스 judge에서는 편향이 더 커지거나 다른 형태로 농축될 수 있다.",
                    "그래서 Phase 3 이후에는 residual position bias와 format failure를 같이 본다.",
                ]),
            ],
            takeaway="원 논문의 진짜 장점은 strong judge를 자랑한 데 있지 않고, judge를 measurement system으로 취급했다는 데 있습니다.",
            notes="""원 논문이 이 정도로 널리 인용된 이유는 결국 숫자가 강했기 때문입니다.
하지만 좋은 paper study는 숫자만 반복하면 안 됩니다.
원 논문은 꽤 인상적인 숫자를 보였지만,
동시에 judge bias도 분석했습니다.

특히 position bias와 verbosity bias를 별도 분석 대상으로 둔다는 점이 중요합니다.
이건 저자들이 judge를 그냥 좋은 채점기로 본 게 아니라,
편향을 가진 measurement system으로 봤다는 뜻입니다.

이 관점이 제 재현 실험과 직접 이어집니다.
저도 agreement 숫자만 반복하지 않고,
Phase 3 이후에는 남는 오류가 어떤 편향 형태로 남는지 따로 분석합니다.

그래서 이 슬라이드는 뒤쪽 residual bias 슬라이드의 이론적 다리 역할을 합니다.""",
        ),
        SlideSpec(
            title="1-7. 그런데 이 논문이 남긴 열린 질문도 있었다",
            section="Base Paper",
            purpose="원 논문을 비판적으로 읽고, 왜 내 재현 실험이 필요한지 자연스럽게 연결한다.",
            bullets=[
                "judge가 강하다는 사실과 judge가 저렴하고 재현 가능하다는 사실은 다르다.",
                "GPT-4는 강력하지만 폐쇄형 API라 비용, 버전 변화, 재현성 문제가 남는다.",
                "bias를 분석했어도 ‘오픈소스 judge도 같은 수준까지 갈 수 있는가’는 아직 열려 있었다.",
            ],
            stat_boxes=[("남은 질문 1", "closed judge"), ("남은 질문 2", "cost / reproducibility"), ("남은 질문 3", "open judge 가능?"), ("내 접근", "faithful reproduction + extension")],
            takeaway="내 실험은 원 논문을 부정하는 게 아니라, 그 프로토콜을 오픈소스 환경으로 옮겼을 때 어디까지 유지되는지 묻는 작업입니다.",
            notes="""좋은 paper study는 찬양으로 끝나면 안 됩니다.
그래서 여기서 원 논문이 남긴 열린 질문을 짚고 넘어가겠습니다.

첫째, GPT-4 judge는 강력하지만 closed judge입니다.
즉 비용과 버전 drift 문제가 남고, 시간이 지나면 같은 실험을 그대로 다시 하기가 어렵습니다.
둘째, 논문은 bias를 분석했지만, 그 결론을 오픈소스 judge에 바로 옮길 수는 없습니다.

제가 이 저장소에서 한 일은 이 열린 질문을 따라가는 것입니다.
MT-Bench 80문항, single/pairwise/reference 프로토콜은 그대로 유지하되,
judge를 Qwen, InternLM, GPT-4o-mini로 나눠서 다시 본 겁니다.

즉 다음 섹션부터는 ‘원 논문을 얼마나 잘 이해했는가’가 아니라
‘그 이해를 바탕으로 무엇을 재현했고 무엇을 추가로 알게 됐는가’로 넘어갑니다.""",
        ),
        SlideSpec(
            title="2부. 내 연구",
            section="Reproduction",
            purpose="앞의 논문 리뷰에서 뒤의 랩미팅 파트로 시선이 전환되는 순간을 명확히 만든다.",
            bullets=[],
            layout="divider",
            takeaway="이제부터는 내 연구입니다. 원 논문의 프로토콜을 오픈소스 judge 실험으로 어떻게 재현하고 어디까지 확장했는지 보여드립니다.",
            notes="""이 슬라이드는 의도적으로 호흡을 한번 끊는 역할을 합니다.
앞 절반은 원 논문 리뷰였고, 여기부터는 제 연구 발표입니다.
그래서 청중이 화면만 봐도 세션이 전환됐다는 걸 느끼게 하고 싶었습니다.

짧게 말하면 앞 절반에서 만든 기준점을 이제 검증 대상으로 바꾸는 순간입니다.
원 논문의 두 핵심 주장, 그리고 열린 질문을 들고 제 실험으로 넘어갑니다.""",
        ),
        SlideSpec(
            title="2-1. 내가 실제로 검증한 네 연구 질문",
            section="Reproduction",
            purpose="뒤 절반의 내 연구 파트를 연구 질문 단위로 묶어 놓고 듣게 한다.",
            bullets=[
                "즉 뒤 절반은 phase log가 아니라 RQ1부터 RQ4까지의 답을 찾아가는 흐름으로 들으면 가장 자연스럽다.",
            ],
            layout="cards",
            cards=[
                ("RQ1. judge scaling", "Qwen judge를 7B → 14B → 32B로 키우면 pairwise reliability가 실제로 단조 개선되는가?"),
                ("RQ2. residual error", "judge가 좋아진 뒤에도 남는 오류는 단순 noise인가, 아니면 순서 민감성과 format failure처럼 구조화되는가?"),
                ("RQ3. ensemble design", "작은 judge 여러 개를 합칠 때, 다수결보다 더 나은 decision rule이 실제로 존재하는가?"),
                ("RQ4. question reduction", "변별도 기반 subset으로 비용을 줄이면서도 모델 서열을 보존할 수 있는가?"),
            ],
            takeaway="이 뒤 절반은 phase 로그가 아니라 네 개의 연구 질문을 차례로 답하는 발표로 듣는 것이 가장 좋습니다.",
            notes="""여기부터는 의도적으로 분위기를 바꿉니다.
divider 슬라이드에서 세션 전환은 이미 끝났습니다.
이 슬라이드는 더 이상 전환을 선언하는 슬라이드가 아니라,
뒤 절반 전체를 묶는 네 연구 질문을 정리하는 슬라이드입니다.

뒤 절반을 들을 때는 네 질문만 기억하면 됩니다.
judge scaling이 실제로 reliability를 개선하는가.
개선 뒤 남는 오류는 어떤 구조인가.
작은 judge들을 합칠 때 어떤 집계 규칙이 더 나은가.
마지막으로 문항 수를 줄여도 서열이 남는가.

즉 이 슬라이드는 랩미팅 파트의 문제 정의 슬라이드라고 보시면 됩니다.""",
        ),
        SlideSpec(
            title="2-2. 실험은 어떤 순서로 진행했는가",
            section="Reproduction",
            purpose="Phase 1부터 Phase 6까지의 실험 순서를 한 장에서 보여준다.",
            bullets=[
                "P1 self-judge 기준선으로 자기평가 편향을 확인하고, P2에서 외부 14B judge로 초기 sanity check를 했다.",
                "P3는 Qwen 7B/14B/32B를 이용한 메인 judge scaling 실험이고, 여기서 RQ1과 RQ2의 핵심 수치가 나온다.",
                "P3의 출력 위에서 RQ3로 ensemble decision rule을 비교하고, P4·P5는 Qwen 결과가 family가 바뀌어도 크게 무너지지 않는지 점검한다.",
                "P6는 11개 모델 풀의 repeated hold-out으로 RQ4, 즉 tinyMT-Bench subset의 운영 가능성을 same-set upper bound와 분리해 다시 본 단계다.",
            ],
            layout="timeline",
            takeaway="즉 오늘 결과는 self-judge 확인 → main scaling → ensemble / cross-family 보강 → repeated hold-out 검증의 순서로 쌓아 올린 것입니다.",
            notes="""이 슬라이드는 청중이 뒤의 결과를 phase log처럼 듣지 않도록 넣은 진행 순서 슬라이드입니다.
실험이 어떤 순서로 쌓였는지 보여주면, 왜 각 phase가 필요한지 더 잘 이해됩니다.

P1은 self-judge bias를 보는 기준선이고, P2는 외부 14B judge로 broad ranking sanity를 확인하는 단계입니다.
이 두 단계는 문제 정의와 예비 검증 역할을 합니다.

핵심은 P3입니다. Qwen 7B, 14B, 32B를 같은 family 안에서 비교해 judge scaling의 메인 결과를 얻습니다.
이 P3 결과 위에서 RQ3인 ensemble design도 평가합니다.
P4와 P5는 그 메인 결과를 보조 검증하는 단계입니다.
InternLM은 cross-family, GPT-4o-mini는 external anchor입니다.

마지막 P6는 question reduction을 같은 모델셋 안에서만 보지 않기 위해 넣은 단계입니다.
즉 Phase 6가 있어야 tinyMT-Bench 결과를 same-set upper bound와 repeated hold-out evidence로 분리해서 말할 수 있습니다.""",
        ),
        SlideSpec(
            title="RQ1 준비. Phase 1–2: 왜 self-judge를 믿으면 안 되는가",
            section="Results",
            purpose="메인 실험에 들어가기 전에 self-judge bias와 초기 sanity check를 짚는다.",
            bullets=[
                "원 논문은 strong judge의 가능성을 보였지만, 이 슬라이드는 오픈소스 self-judge로 가면 바로 어떤 함정이 생기는지 보여줍니다.",
                "Qwen2.5-7B self-judge는 overall 8.12인데 Math·Coding은 각각 8.80으로 튀며 자기 강점을 과대평가한다.",
                "외부 14B judge를 붙이면 broad ranking은 정리되지만, pairwise inconsistency 자체는 여전히 높다.",
                "즉 문제는 ‘judge를 쓰느냐’가 아니라 ‘어떤 judge를 어떻게 쓰느냐’이다.",
            ],
            images=[FIG / "fig0_phase1_scores.png", FIG / "fig2_overall_rankings.png"],
            image_captions=["Phase 1 self-judge category scores", "초기 overall ordering sanity check"],
            takeaway="RQ1을 제대로 보려면 먼저 self-judge 함정을 걷어내야 합니다. 이 슬라이드는 그 준비 단계입니다.",
            notes="""Phase 1과 2는 본격 결과라기보다 문제 제기 단계입니다.
Self-judge를 해보면 Qwen2.5-7B가 특히 Math와 Coding에서 높은 점수를 주는 경향이 드러납니다.
즉 자신의 강점 영역을 과대평가할 가능성이 있습니다.

그래서 외부 judge를 붙이지만, 거기서도 문제가 끝나지 않습니다.
ranking은 어느 정도 정리되지만 pairwise inconsistency가 여전히 높습니다.
이건 judge reliability 문제를 단순히 self-vs-external로 읽으면 안 된다는 뜻입니다.

이 슬라이드의 역할은 뒤 메인 실험의 필요성을 만드는 것입니다.
왜 32B 같은 더 큰 judge를 보고, 왜 cross-family와 hold-out까지 보게 되었는지 설명하는 출발점이라고 생각하시면 됩니다.""",
        ),
        SlideSpec(
            title="RQ1 답. Phase 3: Qwen judge scaling이 메인 결과다",
            section="Results",
            purpose="judge scaling이 reliability에 주는 효과를 가장 강하게 보여준다.",
            bullets=[
                "Pairwise inconsistency는 7B 78.75% → 14B 46.85% → 32B 32.86%로 단조 감소한다.",
                "Single-grade score range도 0.84 → 1.12 → 1.48로 확대되어 모델 간 변별력이 커진다.",
                "따라서 메인 결론은 ‘Qwen 기반 judge scaling의 same-family empirical trend’로 읽는 것이 맞다.",
            ],
            images=[FIG / "fig4_judge_scaling.png", FIG / "fig5_phase3_scores.png"],
            image_captions=["Judge scaling과 카테고리별 inconsistency", "Phase 3 single-grade 점수 분포"],
            stat_boxes=[("7B", "78.75%"), ("14B", "46.85%"), ("32B", "32.86%"), ("핵심", "same-family trend")],
            takeaway="RQ1의 현재 답은 예입니다. judge를 키우면 좋아집니다. 하지만 그 다음에는 RQ2를 물어야 합니다.",
            notes="""이 슬라이드가 이번 재현 실험의 중심입니다.
Qwen judge를 7B에서 14B, 32B로 키우면 pairwise inconsistency가 크게 줄어듭니다.
그리고 score range도 넓어져서 모델 간 차이를 더 잘 구분합니다.

이 결과는 꽤 강합니다. 하지만 표현은 조심해야 합니다.
이건 universal scaling law가 아닙니다.
정확하게는 Qwen2.5 동일 family 안에서 judge가 커질수록 reliability가 좋아졌다는 empirical trend입니다.

저는 발표에서 이 보수적 framing을 일부러 유지할 겁니다.
왜냐하면 뒤의 InternLM과 GPT-4o-mini가 이 메인 결과를 보조하긴 하지만,
Phase 3 자체를 완전히 대체하진 않기 때문입니다.

하지만 이 슬라이드를 그냥 ‘스케일링이 먹혔다’로 끝내면 재미가 없습니다.
바로 다음 슬라이드에서 남는 오류의 구조를 보면,
단순 개선 서사가 아니라 훨씬 더 흥미로운 이야기가 나오기 때문입니다.""",
        ),
        SlideSpec(
            title="RQ2. judge가 좋아진 뒤에도 남는 오류는 무엇인가",
            section="Results",
            purpose="전체 서열과 질문 단위 결정이 분리된다는 점, 그리고 잔여 오류가 순서 민감하게 남는다는 점을 보여준다.",
            bullets=[
                "Qwen32, InternLM20B, GPT-4o-mini는 broad ranking은 대체로 맞지만, exact pairwise winner agreement는 0.50~0.58 수준에 머문다.",
                "특히 Qwen32의 남은 불일치 중 94.93%가 first-position win으로 연결되어, 잔여 오류가 순서 민감한 사례에 집중된다는 점이 드러난다.",
                "즉 ranking validity와 question-level decision cleanliness는 분리해서 읽어야 한다.",
            ],
            images=[FIG / "fig11_position_bias.png", FIG / "fig13_ensemble_v2.png"],
            image_captions=["Order-sensitive residual errors", "Majority vs abstain ensemble"],
            takeaway="RQ2의 답은 명확합니다. judge가 좋아져도 남는 오류는 단순 노이즈가 아니라 순서 민감성과 운영 리스크의 형태로 남습니다.",
            notes="""Phase 4와 5를 보고 나면 많은 청중이 안심합니다.
Qwen32와 InternLM20B, GPT-4o-mini가 대체로 비슷한 서열을 준다면
이제 judge 문제는 거의 해결된 것처럼 느껴질 수 있습니다.

하지만 바로 이 지점에서 한 걸음 더 들어가야 합니다.
32B judge는 전체 inconsistency는 많이 줄였지만,
남아 있는 불일치는 거의 first-position bias와 연결됩니다.

즉 큰 흐름의 ranking validity와, 개별 질문 수준의 decision cleanliness는 다른 문제입니다.
앙상블 결과도 같은 교훈을 줍니다.
작은 judge를 그냥 다수결로 합치면 오히려 오염되고,
abstain ensemble처럼 불확실성을 관리해야 오히려 더 좋아집니다.

그래서 이 슬라이드는 judge 연구를 단순 모델 비교가 아니라
evaluation system design 문제로 읽어야 한다는 메시지를 줍니다.""",
        ),
        SlideSpec(
            title="RQ3. 왜 majority보다 abstain이 더 낫나",
            section="Results",
            purpose="기권 설계를 결과표가 아니라 decision rule 자체로 설명한다.",
            bullets=[
                "각 judge는 pairwise 한 쌍에 대해 {A, B, tie, inconsistent} 중 하나를 남긴다. 여기서 inconsistent는 AB/BA swap이 충돌한 low-confidence case다.",
                "다수결은 inconsistent도 하나의 표처럼 세기 때문에 [inconsistent, inconsistent, A] 같은 경우 winner를 잃고 noisy 7B가 aggregate를 오염시킨다.",
                "abstain은 inconsistent를 기권으로 버리고, 남은 decisive vote가 충돌하지 않을 때만 winner를 선언한다.",
                "실제로 604쌍(36%)이 inconsistent→winner로 복구되고, inconsistency는 58.63%→24.70%, decisive rate는 41.37%→75.30%로 개선된다.",
            ],
            stat_boxes=[("Majority", "58.63%"), ("Abstain", "24.70%"), ("Recovered", "604쌍"), ("Decisive", "75.30%")],
            takeaway="RQ3의 답은 예입니다. 작은 judge를 그냥 다수결로 합치는 것보다, low-confidence vote를 기권으로 다루는 것이 훨씬 더 낫습니다.",
            notes="""이 슬라이드는 제가 이번 발표에서 꼭 분리해서 설명하고 싶은 부분입니다.
많은 청중이 기권 설계를 보면 처음에는 보수적으로 포기한 것처럼 받아들입니다.
하지만 실제로는 반대입니다. 이건 decision rule을 더 정교하게 만든 겁니다.

pairwise 한 쌍에 대해 각 judge는 A, B, tie, inconsistent 중 하나를 냅니다.
여기서 inconsistent는 세 번째 class라기보다, AB/BA swap이 충돌한 low-confidence signal입니다.
그런데 다수결은 이걸 일반 표처럼 셉니다.
그러면 [inconsistent, inconsistent, A] 같은 상황에서도 winner를 잃습니다.

abstain은 다르게 봅니다.
inconsistent를 vote가 아니라 기권으로 두고, 남은 decisive vote가 서로 충돌하지 않을 때만 winner를 선언합니다.
그래서 [inconsistent, inconsistent, A]는 A로 복구되고,
[A, B, inconsistent]는 여전히 inconsistent로 남습니다.

즉 이 설계는 무조건 많이 판단하려는 것도 아니고, 무조건 보수적으로 포기하는 것도 아닙니다.
낮은 품질 judge의 불확실한 표를 집계에서 분리해 measurement noise를 줄이는 방식입니다.
604쌍이 복구되고 decisive rate도 올라간다는 숫자가 바로 그 설계적 의미를 보여줍니다.""",
        ),
        SlideSpec(
            title="RQ4-1. 문항을 줄여도 같은 모델셋에서는 서열이 남는가",
            section="Results",
            purpose="tinyMT-Bench의 same-set upper bound와 random baseline을 함께 보여준다.",
            bullets=[
                "TopDisc-40은 동일 7개 모델 집합에서 ρ=1.000, TopDisc-25는 ρ=0.964를 달성해 same-set upper bound로는 매우 강하다.",
                "반면 random subset은 평균적으로는 좋아져도 분산이 커서, 30문항에서는 mean ρ≈0.95여도 worst-case는 여전히 흔들린다.",
                "즉 same-set 결과만 보면 40문항으로 절반 절감이 가능해 보이지만, 이 수치만으로 운영점을 확정하면 과감해진다.",
            ],
            images=[FIG / "fig9_tiny_mt_bench.png", FIG / "fig7_qsize_sensitivity.png"],
            image_captions=["same-set tinyMT-Bench upper bound", "question count sensitivity of random subsets"],
            stat_boxes=[("TopDisc-25", "ρ=0.964"), ("TopDisc-40", "ρ=1.000"), ("Random 30", "mean ρ≈0.95"), ("해석", "same-set upper bound")],
            takeaway="RQ4의 중간 답은 예입니다. 다만 40문항 결과는 same-set upper bound로 읽고, 운영점은 hold-out에서 다시 확인해야 합니다.",
            notes="""이제 RQ4로 넘어갑니다.
여기서는 먼저 same-set 결과만 따로 보겠습니다.

TopDisc-40은 동일 7개 모델 집합에서는 ρ=1.000입니다.
그래서 얼핏 보면 80문항을 바로 40문항으로 줄여도 되는 것처럼 보입니다.
하지만 바로 옆 random sensitivity를 보면 질문 수가 줄어들수록 분산이 커지고,
mean과 worst-case가 다르다는 점이 드러납니다.

그래서 이 슬라이드는 일부러 strong same-set result를 보여주되,
동시에 이것이 upper bound라는 점을 같이 말하는 슬라이드입니다.
즉 같은 모델셋 안에서는 40문항이 매우 강하지만,
실제 운영점은 다음 슬라이드의 hold-out까지 보고 정해야 합니다.""",
        ),
        SlideSpec(
            title="RQ4-2. hold-out과 다른 judge에서 안전한 운영점은 어디인가",
            section="Results",
            purpose="cross-family judge sanity와 repeated hold-out 운영점을 함께 묶어 same-set 결과를 보수적으로 닫는다.",
            bullets=[
                "InternLM20B는 Qwen32와 ρ=0.893, GPT-4o-mini는 ρ=0.964로 broad ranking pattern을 유지해 Qwen 결과가 완전한 family artifact는 아님을 보여준다.",
                "다만 exact pairwise agreement는 0.50~0.58 수준이라 judge 간 broad consistency와 question-level cleanliness는 여전히 구분해야 한다.",
                "Repeated hold-out 330 split에서는 40문항도 strong하지만, 세 judge 모두에서 가장 안전한 운영점은 60문항(mean ρ≈0.998 / 0.995 / 0.972)이다.",
            ],
            images=[FIG / "fig16_phase345_judge_summary.png", FIG / "fig15_tiny_mt_bench_generalization.png"],
            image_captions=["Cross-family and external judge summary", "330-split repeated hold-out generalization"],
            stat_boxes=[("InternLM20B", "ρ=0.893"), ("GPT-4o-mini", "ρ=0.964"), ("60문항", "three-judge safe zone"), ("해석", "preliminary but strong")],
            takeaway="따라서 RQ4의 최종 답은 ‘부분적으로 예’입니다. 공격적인 40문항은 same-set에서, 더 안전한 운영점은 hold-out 기준 60문항에서 찾는 것이 맞습니다.",
            notes="""이 슬라이드는 RQ4를 보수적으로 닫기 위한 두 번째 슬라이드입니다.

첫 번째 축은 judge family입니다.
InternLM20B와 GPT-4o-mini가 Qwen32와 broad ranking을 꽤 잘 맞춘다는 점은,
Qwen 결과가 완전히 특이한 artifact는 아니라는 근거를 줍니다.
하지만 pairwise exact agreement는 여전히 낮기 때문에,
질문 수준의 cleanliness까지 같다고 말하면 안 됩니다.

두 번째 축은 hold-out입니다.
330개 split 반복 검증을 보면 40문항도 평균적으로 강하지만,
세 judge 모두에서 가장 안전한 운영점은 60문항 쪽입니다.

그래서 발표에서는 40문항을 headline으로 과장하지 않고,
same-set upper bound와 repeated hold-out safe zone을 분리해서 말하는 것이 핵심입니다.""",
        ),
        SlideSpec(
            title="2-3. 네 연구 질문에 대한 현재 답",
            section="Wrap-up",
            purpose="2부에서 처음 던진 네 연구 질문에 대한 현재 답을 한 장에서 닫아 준다.",
            bullets=[
                "즉 이 발표의 결론은 ‘오픈소스 judge도 충분하다’가 아니라, strong judge 메시지의 경계와 운영 규칙까지 같이 말할 수 있게 되었다는 데 있습니다.",
            ],
            layout="cards",
            cards=[
                ("RQ1. judge scaling은 개선되는가", "예. Qwen judge scaling 결과에서는 7B → 14B → 32B로 갈수록 pairwise inconsistency가 뚜렷하게 줄었습니다."),
                ("RQ2. 남는 오류는 어떤 구조인가", "단순 noise가 아니라 순서 민감성과 format failure처럼 구조화됩니다. 그래서 ranking과 question-level cleanliness를 분리해야 합니다."),
                ("RQ3. 어떤 앙상블이 더 낫나", "단순 majority보다 abstain이 낫습니다. inconsistent를 표로 세지 않을 때 더 많은 쌍을 더 깨끗하게 결정할 수 있습니다."),
                ("RQ4. 문항을 줄여도 되나", "부분적으로 예. same-set upper bound는 40문항에서 강하지만, repeated hold-out 기준 안전한 운영점은 60문항 쪽입니다."),
            ],
            takeaway="따라서 결론은 단순 재현 성공이 아니라, 오픈소스 judge를 어디까지 믿고 어디서부터 운영 규칙을 붙여야 하는지 실험적으로 말할 수 있게 되었다는 것입니다.",
            notes="""이 슬라이드는 2부 전체를 닫아 주는 슬라이드입니다.
2부가 phase 보고의 나열로 끝나면 연구 질문이 희미해집니다.
그래서 처음 던진 RQ1부터 RQ4까지를 다시 불러와 답하는 방식으로 마무리합니다.

RQ1에 대한 답은 가장 명확합니다.
Qwen judge scaling 결과에서는 reliability 개선이 분명히 관찰됐습니다.
RQ2에 대한 답도 선명합니다.
남는 오류는 단순 noise가 아니라 순서 민감성과 format failure입니다.

RQ3는 decision rule 문제였습니다.
작은 judge를 그냥 다수결로 세면 안 되고, low-confidence vote를 기권으로 다뤄야 한다는 점이 드러났습니다.

RQ4는 부분적으로만 예라고 답해야 정직합니다.
40문항 headline만 남기면 과하고, same-set upper bound와 hold-out safe zone을 같이 말해야 합니다.

이 슬라이드가 들어가면 2부가 phase 보고가 아니라
네 개의 연구 질문에 대한 실험적 답변으로 닫히게 됩니다.""",
        ),
        SlideSpec(
            title="질문 전에: 저장소는 이렇게 읽으면 됩니다",
            section="Q&A",
            purpose="발표 후 repo를 실제로 볼 사람을 위해 마지막 안내를 남기고 Q&A로 넘긴다.",
            bullets=[
                "원 논문 설명은 `presentation/`의 deck와 notes, 재현 서사는 `README.md`, 원고는 `paper/`를 보면 된다.",
                "실험은 Phase 1–2 → Phase 3 → Phase 4–5 → Phase 6 순으로 읽으면 오늘 발표 흐름과 정확히 맞는다.",
                "수치 검증은 `figures/` 옆 CSV와 `data/` raw judgment를 따라가면 된다.",
                "즉 오늘 발표는 끝나도, repo 안에서 결론을 다시 확인할 수 있게 설계되어 있다.",
            ],
            stat_boxes=[("presentation", "today's deck"), ("README", "story"), ("data", "evidence"), ("paper", "final text")],
            takeaway="질문 받겠습니다.",
            notes="""마지막 슬라이드는 발표를 마무리하면서 동시에 repo 안내를 하는 슬라이드입니다.
교수님이 발표 후 저장소를 실제로 보신다면 어디부터 보면 되는지를 남기는 역할도 합니다.

오늘 발표는 paper study였기 때문에, 발표가 끝나면 질문이 두 방향으로 나올 겁니다.
원 논문 자체에 대한 질문과, 제 재현 실험에 대한 질문입니다.
전자는 presentation과 paper를 보면 되고,
후자는 README, figures, data를 따라가면 됩니다.
특히 저장소는 phase 순서로 읽으면 오늘 발표의 2부와 같은 흐름이 다시 재생됩니다.

이상으로 발표를 마치고 질문 받겠습니다.""",
        ),
    ]


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides = get_slide_specs()
    total = len(slides)
    for idx, spec in enumerate(slides, 1):
        if idx == 1:
            build_title_slide(prs, spec, total)
        elif spec.layout == "divider":
            build_divider_slide(prs, spec, idx, total)
        elif spec.layout == "overview":
            build_overview_slide(prs, spec, idx, total)
        elif spec.layout == "cards":
            build_cards_slide(prs, spec, idx, total)
        elif spec.layout == "compare":
            build_compare_slide(prs, spec, idx, total)
        elif spec.layout == "timeline":
            build_timeline_slide(prs, spec, idx, total)
        else:
            build_content_slide(prs, spec, idx, total)
    prs.save(PPTX_PATH)
    build_notes_md(slides)
    print(PPTX_PATH)
    print(NOTES_MD_PATH)


if __name__ == "__main__":
    build_deck()
